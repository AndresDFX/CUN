# -*- coding: utf-8 -*-
"""
Motor reutilizable para generar diapositivas .pptx con la identidad REAL de la CUN (Corporación
Unificada Nacional de Educación Superior), tomada de https://cun.edu.co/ (logo real descargado +
colores extraídos del CSS de producción del sitio — ver config/universidades/cun.json._fuente).

PALETA REAL (extraída 2026-08-07 de https://cun.edu.co/ CSS en vivo + logo oficial):
  Verde CUN #007433 (marca primaria) · Verde lima #91DC00 (acento puntual) ·
  Azul marino #0C2340 (encabezados/bandas) · Naranja acento #FF9E1C · Tipografía Calibri.
  Logo real: config/slides/assets/logo_cun_solo.png (wordmark limpio, sin sello "40 años").

Incluye `fechas_inicio_fin_slide()`: cronograma en tarjetas claras de INICIO / CIERRE por bloque
(ACA), más legible que una barra proporcional. `timeline_slide()` se conserva por compatibilidad.

USO:
    from cun_slides_engine import *
    prs = new_prs()
    course_cover(prs, materia, subtitulo, meta_lines)  # grupo(s) solo aquí (portada)
    tutor_slide(prs, "Docente", credenciales, correo, idx=2)  # título genérico; sin nombre propio
    content_slide(prs, "Título", ["– viñeta", ("● subviñeta", 1)], idx=3)
    contenido_sesiones_slide(prs, sessions, idx=4)  # CONTENIDO · Sesión N — tema
    table_content(prs, "Título", ["A","B"], [["x","y"]], idx=5)
    timeline_slide(prs, "Título", date(2026,8,3), date(2026,11,22), blocks, marks, dots, idx=5)
    box_note_slide(prs, "Título", [("aclaracion","Texto…"), ("advertencia","Texto…")], idx=6)
    closing_slide(prs, "Mensaje", ["línea"], "Acento")
    prs.save("salida.pptx")

Pie de página: **vacío** en los cursos CUN (los builds llaman `set_footer("")`). Solo se
imprime el nº de slide a la derecha. La hora de **inicio efectivo** (horario oficial + 10 min)
aparece **una sola vez**, en la portada de la Presentación del Curso, vía `cover_meta_lines`
de `carga_academica` — decisión del docente del 2026-08-10 (antes se repetía en cada slide).

Requiere: pip install python-pptx pillow
"""
import csv, os, re, sys, datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Paleta CUN real (ver cun.json — CSS vivo cun.edu.co) ----------
VERDE  = RGBColor(0x00, 0x74, 0x33)   # marca primaria #007433 (CSS del sitio)
LIMA   = RGBColor(0x91, 0xDC, 0x00)   # acento puntual (líneas/detalles pequeños, NUNCA fondos grandes)
NAVY   = RGBColor(0x0C, 0x23, 0x40)   # azul marino — encabezados, bandas de título, tablas
SLATE  = RGBColor(0x32, 0x4A, 0x6D)   # azul secundario
NARANJA_ISO = RGBColor(0xFF, 0x9E, 0x1C)  # naranja de acento del sitio (#FF9E1C)
INFO   = RGBColor(0xE8, 0xF5, 0xEC)   # caja info / principio institucional (verde muy claro)
ACLAR  = RGBColor(0xFD, 0xEC, 0xD8)   # caja de aclaración metodológica
WARN   = RGBColor(0xFB, 0xE4, 0xE4)   # caja de advertencia crítica
GRAY   = RGBColor(0x2B, 0x2B, 0x2B)   # texto principal
BANNER = NAVY
ALT    = RGBColor(0xF2, 0xF2, 0xF2)   # filas alternas de tabla
SOFT   = RGBColor(0x8F, 0x98, 0x9D)   # gris medio (número de slide)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xC0, 0x39, 0x2B)   # festivos en la línea de tiempo
FONT   = "Calibri"

EMU_IN = 914400
SW, SH = 13.333, 7.5
MARGIN = 0.7
CONTENT_W = SW - 2 * MARGIN
FONT_SCALE = 1.0
_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
LOGO = os.path.join(_ASSETS, "logo_cun_solo.png")   # wordmark limpio (sin sello "40 años")
LOGO_ANIV = os.path.join(_ASSETS, "logo_cun.png")   # versión con sello de aniversario, si se necesita
GMAIL = os.path.join(_ASSETS, "gmail_icon.png")
QR_PRESENTACION_ESTUDIANTES = os.path.join(_ASSETS, "qr_presentacion_estudiantes.png")
# Tablero rompehielos / Preséntate. YA NO es «el mismo en los 5 cursos»: desde el
# 2026-08-11 solo lo usan los grupos pequeños (ver `modo_rompehielos`).
PADLET_PRESENTACION_URL = "https://padlet.com/andres_dfx/cun-wruz81hmf9k06gd7"
LINK_TEAL = RGBColor(0x00, 0x74, 0x33)  # enlaces en verde de marca CUN

# Pie de slides: hora de inicio efectivo (horario + 10 min). Sin nombre del curso.
_FOOTER = ""

def set_footer(text=None):
    """Define el pie de texto de las slides siguientes.

    **Decisión 2026-08-10 (docente):** en los cursos CUN el pie va **vacío**. La hora de
    inicio efectivo ya NO se repite en todas las slides: aparece una sola vez, en la
    **portada de la Presentación del Curso** (vía `cover_meta_lines`). Los builds CUN
    llaman `set_footer("")`; el número de slide sigue saliendo a la derecha.
    """
    global _FOOTER
    _FOOTER = (text or "").strip()

# ---------- Texto enriquecido: **negrita** y @@resaltado en VERDE de marca@@ ----------
def _run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    f = run.font
    f.size = Pt(round(size * FONT_SCALE, 1)); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color

def _rich(paragraph, markup, size, base_color, bold=False, italic=False):
    for part in re.split(r'(@@.*?@@|\*\*.*?\*\*)', markup):
        if part == "":
            continue
        run = paragraph.add_run()
        if part.startswith("@@") and part.endswith("@@"):
            _run(run, part[2:-2], size, VERDE, bold=True, italic=italic)
        elif part.startswith("**") and part.endswith("**"):
            _run(run, part[2:-2], size, base_color, bold=True, italic=italic)
        else:
            _run(run, part, size, base_color, bold=bold, italic=italic)

# ---------- Primitivas ----------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg_white(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = WHITE

def rect(slide, x, y, w, h, color, line=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if not line:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def rounded(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def _logo_hw(width):
    """Alto en pulgadas del logo para un ancho dado (respeta el aspect ratio real)."""
    try:
        from PIL import Image as _PILImage
        iw, ih = _PILImage.open(LOGO).size
    except Exception:
        iw, ih = (150, 104)
    return width * ih / iw

def add_logo(slide, width=1.7, mr=0.5, mt=0.35, corner="right-top"):
    """Logo real de la CUN. corner: 'right-top' | 'left-top' | 'right-bottom' | 'left-bottom'.

    Devuelve el borde inferior (en pulgadas) para que el llamador coloque contenido
    debajo sin superponerse — ver `course_cover`.
    """
    if not os.path.exists(LOGO):
        return mt
    hi = _logo_hw(width)
    w = Inches(width); h = Inches(hi)
    left = Inches(mr) if "left" in corner else Emu(int(SW * EMU_IN) - w - Inches(mr))
    top = Inches(mt) if "top" in corner else Emu(int(SH * EMU_IN) - h - Inches(mt))
    slide.shapes.add_picture(LOGO, left, top, width=w)
    return (mt + hi) if "top" in corner else (SH - mt - hi)

def _fit_title_size(text, width_in, max_lines=2, base=34, min_size=20):
    """Cuerpo de fuente que hace caber `text` en `max_lines` dentro de `width_in`.

    Estimación tipográfica (Calibri bold ≈ 0.5·em de ancho medio por carácter);
    evita que un título largo desborde la banda o invada la zona del logo.
    """
    import math as _math
    text = text or ""
    size = base
    while size > min_size:
        cpl = max(1, int(width_in / (size * 0.5 / 72.0)))
        if _math.ceil(len(text) / cpl) <= max_lines:
            return size
        size -= 2
    return min_size

def footer_num(slide, idx):
    """Pie: hora de inicio efectivo (izq.) + nº de slide (der.). Sin nombre del curso."""
    if _FOOTER:
        # Tras zona de logo izq. (tutor_slide / portadas) para no solaparse.
        tf = textbox(slide, 2.4, SH - 0.45, 6.5, 0.3)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _run(p.add_run(), _FOOTER, 10, SOFT, bold=True)
    if idx is not None:
        tf = textbox(slide, SW - 1.2, SH - 0.45, 0.6, 0.3)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        _run(p.add_run(), str(idx), 10, SOFT, bold=True)

def title_block(slide, title, sub=None):
    """Banda AZUL MARINO con el título centrado en blanco (identidad CUN real)."""
    rect(slide, 0, 0, SW, 1.3, BANNER)
    tf = textbox(slide, MARGIN, 0, CONTENT_W - 2.0, 1.3, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p.add_run(), title, 26, WHITE, bold=True)
    if sub:
        ts = textbox(slide, MARGIN, 1.4, CONTENT_W, 0.45)
        ps = ts.paragraphs[0]; ps.alignment = PP_ALIGN.CENTER
        _rich(ps, sub, 15, GRAY)
    return 2.05 if sub else 1.7

def bullets(slide, items, top, size=16, width=None, left=None):
    left = MARGIN if left is None else left
    width = CONTENT_W if width is None else width
    tf = textbox(slide, left, top, width, SH - top - 0.6)
    for i, it in enumerate(items):
        text, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(10 if lvl == 0 else 6)
        marker = "–   " if lvl == 0 else "●   "
        _rich(p, marker + text, size - (1 if lvl else 0), GRAY)
    return tf

# ---------- Tablas ----------
def _nogrid(table):
    tblPr = table._tbl.tblPr
    for c in list(tblPr):
        if c.tag == qn('a:tableStyleId'):
            tblPr.remove(c)
    sid = tblPr.makeelement(qn('a:tableStyleId'), {})
    sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'  # No Style, No Grid
    tblPr.append(sid)
    table.first_row = False; table.horz_banding = False

def _fill(cell, color):
    cell.fill.solid(); cell.fill.fore_color.rgb = color

def _cell(cell, markup, size=12, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    _rich(p, markup, size, color, bold=bold)

def table_raw(slide, headers, rows, top, col_w=None, header_color=NAVY,
              fs_head=13, fs_body=12, aligns=None):
    ncol = len(headers); nrow = len(rows) + 1
    col_w = col_w or [CONTENT_W / ncol] * ncol
    s = sum(col_w); col_w = [w / s * CONTENT_W for w in col_w]
    height = min(SH - top - 0.6, 0.55 * nrow)
    table = slide.shapes.add_table(nrow, ncol, Inches(MARGIN), Inches(top),
                                   Inches(CONTENT_W), Inches(height)).table
    _nogrid(table)
    for j, w in enumerate(col_w):
        table.columns[j].width = Inches(w)
    aligns = aligns or [PP_ALIGN.LEFT] * ncol
    for j, h in enumerate(headers):
        c = table.cell(0, j); _fill(c, header_color)
        _cell(c, h, size=fs_head, color=WHITE, bold=True, align=aligns[j])
    for i, row in enumerate(rows):
        rc = WHITE if i % 2 == 0 else ALT
        for j, val in enumerate(row):
            c = table.cell(i + 1, j); _fill(c, rc)
            _cell(c, val, size=fs_body, color=GRAY, align=aligns[j])
    return table, top + height

# ---------- Slides de contenido ----------
def content_slide(prs, title, items, sub=None, size=16, idx=None, **_):
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    bullets(s, items, top=top + 0.15, size=size)
    footer_num(s, idx)
    return s


def contenido_sesiones_slide(prs, sessions, idx=None, title="CONTENIDO",
                             max_per_slide=None, size=None):
    """Listado de sesiones en UNA sola slide (Presentación del Curso).

    Título en banda (CONTENIDO) + viñetas ``**Sesión N** — tema.``
    Sin fechas, sin tablas de cronograma. Nunca parte el listado en 2+ slides.
    Si hay muchas sesiones (p. ej. P1 11, TG2 12, TG3 16), usa tipografía
    compacta y, a partir de ~10 ítems, **dos columnas en la misma slide**.

    ``sessions``: iterable de dicts ``{n, titulo}`` o tuplas ``(n, titulo)``.
    ``max_per_slide``: ignorado (compat API; siempre 1 slide).
    ``idx``: número de la slide.
    Returns: siempre ``1`` (cantidad de slides creadas).
    """
    del max_per_slide  # compat: ya no se parte el listado
    items_all = []
    for s in sessions or []:
        if isinstance(s, dict):
            n = s.get("n") if s.get("n") is not None else s.get("num")
            titulo = (s.get("titulo") or s.get("tema") or "").strip()
        else:
            n, titulo = s[0], str(s[1]).strip()
        if not titulo and n is None:
            continue
        if titulo.endswith("."):
            titulo = titulo[:-1].rstrip()
        items_all.append(f"**Sesión {int(n)}** — {titulo}.")

    if not items_all:
        content_slide(prs, title, ["(Sin sesiones definidas.)"], idx=idx, size=size or 15)
        return 1

    n_total = len(items_all)
    use_two_cols = n_total >= 10

    if size is not None:
        fs = size
    elif n_total <= 7:
        fs = 15
    elif n_total <= 9:
        fs = 13
    elif n_total <= 12:
        fs = 12
    else:
        fs = 11

    s = blank(prs)
    bg_white(s)
    top = title_block(s, title)
    y = top + 0.08
    space = Pt(4 if n_total >= 12 else (6 if n_total >= 10 else 8))

    if not use_two_cols:
        tf = textbox(s, MARGIN, y, CONTENT_W, SH - y - 0.55)
        for i, text in enumerate(items_all):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = 0
            p.space_after = space
            _rich(p, "–   " + text, fs, GRAY)
    else:
        mid = (n_total + 1) // 2
        col_w = (CONTENT_W - 0.25) / 2
        gap = 0.25
        for col_i, chunk in enumerate((items_all[:mid], items_all[mid:])):
            left = MARGIN + col_i * (col_w + gap)
            tf = textbox(s, left, y, col_w, SH - y - 0.55)
            for i, text in enumerate(chunk):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.level = 0
                p.space_after = space
                _rich(p, "–   " + text, fs, GRAY)

    footer_num(s, idx)
    return 1


def table_content(prs, title, headers, rows, sub=None, note=None, col_w=None,
                  aligns=None, idx=None, fs_body=12, **_):
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    _, bottom = table_raw(s, headers, rows, top=top + 0.15, col_w=col_w, aligns=aligns, fs_body=fs_body)
    if note:
        _rich(textbox(s, MARGIN, min(SH - 0.85, bottom + 0.15), CONTENT_W, 0.6).paragraphs[0],
              note, 12, GRAY, italic=True)
    footer_num(s, idx)
    return s

# ---------- Cajas de color (aclaración / advertencia / info) ----------
_BOX_COLORS = {"aclaracion": ACLAR, "advertencia": WARN, "info": INFO}
_BOX_LABELS = {"aclaracion": "💡 Aclaración", "advertencia": "⚠️ Importante", "info": "ℹ️ Nota institucional"}

def box_note_slide(prs, title, boxes, sub=None, idx=None):
    """boxes: lista de tuplas (tipo, texto) con tipo en {'aclaracion','advertencia','info'}."""
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    y = top + 0.2
    for tipo, texto in boxes:
        color = _BOX_COLORS.get(tipo, INFO)
        label = _BOX_LABELS.get(tipo, "")
        h = 0.5 + 0.28 * max(1, (len(texto) // 90) + 1)
        rounded(s, MARGIN, y, CONTENT_W, h, color)
        tf = textbox(s, MARGIN + 0.25, y + 0.12, CONTENT_W - 0.5, h - 0.24)
        p = tf.paragraphs[0]; _rich(p, "**" + label + "**", 13, NAVY, bold=True)
        p2 = tf.add_paragraph(); p2.space_before = Pt(4); _rich(p2, texto, 13, GRAY)
        y += h + 0.25
    footer_num(s, idx)
    return s

def tutor_slide(prs, nombre, credenciales, correo, rol=None, idx=None):
    """Diapositiva del docente: título genérico «Docente», bullets de perfil y correo.

    No proyecta el nombre propio (el param `nombre` se ignora para el título; compat API).
    Perfil canónico CUN (pasa la lista desde el build): Ingeniero de Sistemas ·
    Candidato a MSc en Inteligencia Artificial · Líder Técnico · Speaker Tecnológico.
    Sin subtítulo tipo «Docente · <curso>» (el param `rol` se ignora; compat API).
    Logo CUN abajo-izquierda. Pie = hora de inicio efectivo (sin nombre del curso).
    """
    s = blank(prs); bg_white(s)
    title_block(s, "Docente")  # genérico; sin nombre propio en pantalla
    tx = 2.8
    tc = textbox(s, tx, 2.35, SW - tx - MARGIN, 2.05)
    for i, ln in enumerate(credenciales):
        p = tc.paragraphs[0] if i == 0 else tc.add_paragraph()
        p.space_after = Pt(12); _rich(p, ln, 18, GRAY)
    ey = 4.55
    if os.path.exists(GMAIL):
        s.shapes.add_picture(GMAIL, Inches(tx), Inches(ey), height=Inches(0.32))
        te = textbox(s, tx + 0.55, ey - 0.02, SW - tx - 0.55 - MARGIN, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    else:
        te = textbox(s, tx, ey - 0.02, SW - tx - MARGIN, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    _run(te.paragraphs[0].add_run(), correo, 16, LINK_TEAL, bold=True)
    add_logo(s, width=1.6, corner="left-bottom", mt=0.35, mr=0.5)
    footer_num(s, idx)
    return s


# ---------------------------------------------------------------------------
# ROMPEHIELOS «PRESÉNTATE» — la forma la decide el TAMAÑO del grupo
# ---------------------------------------------------------------------------
# Decisión del docente (2026-08-11), con la matrícula real auditada en CDigital:
#   Investigación 53339 = 20 · Proyecto I 54ES4 = 50 · Creatividad 54408 = 50 ·
#   TG2 54448 = 50 · TG3 (54450 + 54466 + 54467) = 112 en UNA sola serie.
#
#   ≤ ICEBREAKER_MAX_MURO  → MURO de Padlet. Un muro de 20 notas se lee entero y todos
#                            alcanzan a ser vistos, que es el punto del rompehielos.
#   >  ICEBREAKER_MAX_MURO → FORMULARIO de Google (gratis, sin tope de participantes, ya
#                            incluido en la licencia CUN: el estudiante entra con su
#                            @cun.edu.co y no crea cuenta) + **encuestas y Q&A nativos de
#                            Meet** para la parte en vivo. Con 50 —o con los 112 de TG3—
#                            el muro no se alcanza a leer, y los planes gratis de
#                            Mentimeter (50/mes) y Slido (100 · 3 encuestas) se quedan
#                            justo por debajo de estos cursos.
#
# La elección NO se escribe curso por curso en los builds: `modo_rompehielos()` la deriva
# de la matrícula, y `contar_estudiantes()` la cuenta de los roster descargados de
# CDigital. Un curso sin roster completo cae al modo grande, que es el que no se rompe.
ICEBREAKER_MAX_MURO = 20
MODO_MURO = "muro"              # Padlet
MODO_FORMULARIO = "formulario"  # Google Forms + encuestas/Q&A de Meet
# Roster: `<carpeta de la asignatura>/<año>/<grupo>/Listado estudiantes (CDigital).csv`.
# El año sale del `inicio` del curso; ROSTER_ANIO es solo el respaldo si el JSON no lo trae.
ROSTER_ANIO = "2026"
ROSTER_CSV = "Listado estudiantes (CDigital).csv"
# Clave del enlace real del formulario en `config/cursos/carga_academica_2026.json`
# → cursos.<key>.formulario_presentacion. Mismo contrato que `meet` y que `clases`:
# fuente única, cadena vacía (o clave ausente) ⇒ el material muestra el placeholder.
FORMULARIO_PRESENTACION_KEY = "formulario_presentacion"

_CARGA_MOD = None
_AVISADO_SIN_CURSO = False


def _carga_academica():
    """Módulo `config/cursos/carga_academica` (import perezoso). None si no está.

    El motor no depende de la carga académica para dibujar: solo la consulta para
    saber cuántos estudiantes tiene el curso y de dónde sale el enlace del formulario.
    """
    global _CARGA_MOD
    if _CARGA_MOD is None:
        import importlib
        ruta = os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "cursos"))
        if ruta not in sys.path:
            sys.path.insert(0, ruta)
        try:
            _CARGA_MOD = importlib.import_module("carga_academica")
        except Exception:  # pragma: no cover — sin config el motor sigue dibujando
            _CARGA_MOD = False
    return _CARGA_MOD or None


def _estudiantes_en_roster(path):
    """Filas de estudiante del roster de CDigital. Devuelve -1 si el archivo no está.

    El CSV descargado del aula trae `nombre,correo,rol`; se cuentan solo las filas con
    rol de estudiante (el Profesor viene en la misma lista). Si el archivo no tiene
    columna `rol` se cuenta toda fila con datos: es una aproximación por exceso, y
    contar de más solo puede llevar al modo formulario, que es el que aguanta.
    """
    if not os.path.isfile(path):
        return -1
    n = 0
    try:
        with open(path, encoding="utf-8-sig", newline="") as f:
            for fila in csv.DictReader(f):
                campos = {(k or "").strip().lower(): (v or "").strip() for k, v in fila.items()}
                rol = campos.get("rol", "")
                if rol:
                    if rol.lower().startswith("estudiante"):
                        n += 1
                elif campos.get("correo") or campos.get("nombre"):
                    n += 1
    except OSError:
        return -1
    return n


def contar_estudiantes(course_key):
    """Matrícula del curso = suma de los roster de CDigital de TODOS sus grupos.

    Ruta de cada roster: ``<carpeta de la asignatura>/2026/<grupo>/Listado estudiantes
    (CDigital).csv`` (los grupos salen de `carga_academica_2026.json`). TG3 suma sus
    tres grupos porque corren en **una sola serie** de encuentros: el rompehielos lo
    hacen juntos, así que el tamaño que importa es el de la serie.

    Devuelve ``None`` cuando no se puede afirmar la matrícula completa —sin config, sin
    grupos, o con el roster de algún grupo ausente—. Quien decide el modo trata ese
    ``None`` como grupo grande.
    """
    carga = _carga_academica()
    if carga is None or not course_key:
        return None
    try:
        c = carga.curso(course_key)
        grupos = list(c.get("groups") or [])
        base = str(carga.course_dir(course_key))
        anio = str(c.get("inicio") or "")[:4] or ROSTER_ANIO
    except Exception:
        return None
    if not grupos:
        return None
    total = 0
    for g in grupos:
        n = _estudiantes_en_roster(os.path.join(base, anio, str(g), ROSTER_CSV))
        if n < 0 and anio != ROSTER_ANIO:
            n = _estudiantes_en_roster(os.path.join(base, ROSTER_ANIO, str(g), ROSTER_CSV))
        if n < 0:
            return None
        total += n
    return total


def modo_rompehielos(course_key=None, n_estudiantes=None):
    """`MODO_MURO` (Padlet) o `MODO_FORMULARIO` (Google Forms + Meet), según el tamaño.

    `n_estudiantes` solo para forzar el cálculo en pruebas; en los builds se deja que lo
    cuente `contar_estudiantes`. Sin matrícula conocida → formulario.
    """
    n = n_estudiantes if n_estudiantes is not None else contar_estudiantes(course_key)
    if n is None:
        return MODO_FORMULARIO
    return MODO_MURO if n <= ICEBREAKER_MAX_MURO else MODO_FORMULARIO


def usa_padlet(course_key):
    """True si a ese curso le toca el muro de Padlet por tamaño de grupo."""
    return modo_rompehielos(course_key) == MODO_MURO


def formulario_presentacion_placeholder(curso_corto):
    return f"[URL Formulario Preséntate — pendiente · {curso_corto}]"


def formulario_presentacion_url(course_key, curso_corto=None):
    """Enlace del formulario «Preséntate» del curso: el real si está en config, si no el
    marcador de posición.

    Mismo contrato que el Meet cuando falta la sala: la URL vive en
    `carga_academica_2026.json` → cursos.<key>.formulario_presentacion y **no** se
    escribe en los builds. Mientras el docente no cree el formulario, la clave está
    vacía (o no existe) y en la slide se ve el placeholder, que es el aviso de que
    falta. Uno por asignatura, para que la hoja de respuestas no mezcle cursos.
    """
    url = ""
    corto = curso_corto or ""
    carga = _carga_academica()
    if carga is not None and course_key:
        try:
            c = carga.curso(course_key)
            url = (c.get(FORMULARIO_PRESENTACION_KEY) or "").strip()
            corto = curso_corto or (c.get("titulo_corto") or course_key)
        except Exception:
            url = ""
    return url or formulario_presentacion_placeholder(corto or course_key or "el curso")


def _qr_formulario(course_key):
    """QR del formulario: el del curso si existe en assets, si no el genérico. '' si no hay.

    No se genera aquí: mientras el formulario no tenga enlace real no hay nada que
    codificar, y la slide muestra el marcador de posición del QR.
    """
    nombres = []
    if course_key:
        nombres.append(f"qr_formulario_presentacion_{course_key}.png")
    nombres.append("qr_formulario_presentacion.png")
    for nombre in nombres:
        p = os.path.join(_ASSETS, nombre)
        if os.path.exists(p):
            return p
    return ""


def _icebreaker_render(prs, idx, sub, items, url, qr, size=14, qr_pendiente=None):
    """Dibuja la slide del rompehielos (bullets a la izquierda + QR grande a la derecha).

    `qr_pendiente`: si el QR todavía no existe **y se espera que no exista** (formulario
    sin enlace), el texto que explica el pendiente; se pinta una tarjeta «QR pendiente»
    en el mismo hueco. Sin ese texto, un QR ausente es un error de assets y se avisa.
    """
    s = blank(prs)
    bg_white(s)
    top = title_block(s, "PRESÉNTATE — ROMPEHIELOS", sub)
    bullets(s, items, top=top + 0.15, left=MARGIN, width=6.4, size=size)
    qr_size = 3.4
    qr_x = SW - MARGIN - qr_size
    qr_y = top + 0.25
    if qr and os.path.exists(qr):
        rounded(s, qr_x - 0.12, qr_y - 0.12, qr_size + 0.24, qr_size + 0.85, ALT)
        s.shapes.add_picture(qr, Inches(qr_x), Inches(qr_y), width=Inches(qr_size), height=Inches(qr_size))
        cap = textbox(s, qr_x - 0.12, qr_y + qr_size + 0.02, qr_size + 0.24, 0.28)
        p = cap.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _rich(p, "**Escanea o abre**", 12, NAVY, bold=True)
        url_box = textbox(s, qr_x - 0.12, qr_y + qr_size + 0.28, qr_size + 0.24, 0.45)
        up = url_box.paragraphs[0]
        up.alignment = PP_ALIGN.CENTER
        _run(up.add_run(), url, 9, VERDE, bold=True)
    elif qr_pendiente:
        rounded(s, qr_x - 0.12, qr_y - 0.12, qr_size + 0.24, qr_size + 0.85, ALT)
        tf = textbox(s, qr_x + 0.15, qr_y + 0.9, qr_size - 0.3, 1.6, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), "QR pendiente", 18, NAVY, bold=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)
        _rich(p2, qr_pendiente, 12, GRAY)
        url_box = textbox(s, qr_x - 0.12, qr_y + qr_size + 0.02, qr_size + 0.24, 0.7)
        up = url_box.paragraphs[0]
        up.alignment = PP_ALIGN.CENTER
        _run(up.add_run(), url, 9, VERDE, bold=True)
    else:
        tf = textbox(s, qr_x, qr_y, qr_size, 1.2)
        _rich(tf.paragraphs[0], "⚠️ QR no encontrado en assets.", 14, RED, bold=True)
    footer_num(s, idx)
    return s


def icebreaker_qr_slide(prs, idx=None, consignas=None, sub=None, qr_path=None,
                        padlet_url=None, course_key=None, modo=None,
                        n_estudiantes=None, pide=None, formulario_url=None):
    """Rompehielos «Preséntate» de la Presentación del Curso (= momento de Sesión 01).

    **Pase siempre `course_key`.** Con él, la slide que sale la decide el tamaño real del
    grupo (ver `modo_rompehielos`): muro de Padlet hasta `ICEBREAKER_MAX_MURO`
    estudiantes, formulario de Google por encima. Ningún build elige el modo a mano.

    - `pide`: lo único que cambia entre cursos — qué se le pide al estudiante además del
      nombre («**estado actual** del proyecto (1 frase) + expectativa de TG2»).
    - `consignas`: reemplaza los bullets completos (escotilla de escape; si la usa, el
      texto deja de seguir al modo y le toca mantenerlo a mano).
    - `modo` / `n_estudiantes`: forzar la decisión (pruebas y casos puntuales).

    Sin `course_key` no hay tamaño que consultar y se sirve el muro, que es lo que hacía
    esta función antes de 2026-08-11; se avisa por consola una sola vez.
    """
    global _AVISADO_SIN_CURSO
    n = n_estudiantes
    if n is None and course_key:
        n = contar_estudiantes(course_key)
    if modo is None:
        if not course_key and n is None:
            if not _AVISADO_SIN_CURSO:
                _AVISADO_SIN_CURSO = True
                print(
                    "[cun_slides_engine] AVISO: icebreaker_qr_slide sin `course_key`; se "
                    "sirve el muro de Padlet por compatibilidad. Pase course_key para que "
                    "la forma del rompehielos la decida el tamaño del grupo.",
                    file=sys.stderr,
                )
            modo = MODO_MURO
        else:
            modo = modo_rompehielos(course_key, n)

    if modo == MODO_MURO:
        url = padlet_url or PADLET_PRESENTACION_URL
        pide_txt = pide or "**expectativa del curso** o **tema de interés** (1 frase)"
        # El número del grupo va en las viñetas, no también en el subtítulo.
        sub_txt = sub or "Tablero colaborativo (Padlet)"
        items = consignas or [
            f"**Escanea el QR o abre:** {url}",
            f"En un post-it escribe: **tu nombre** + {pide_txt}.",
            (f"**Ahora, ~7 min.** Somos {n}: el muro se lee entero y nadie queda sin ser visto."
             if n is not None else "**Ahora, ~7 min.** El muro se lee entero: nadie queda sin ser visto."),
            "Después las leemos en voz alta y agrupamos expectativas — sin juzgar: de ahí salen "
            "los ejemplos con los que trabajamos el resto del curso.",
        ]
        return _icebreaker_render(prs, idx, sub_txt, items, url,
                                  qr_path or QR_PRESENTACION_ESTUDIANTES, size=14)

    # ---- Modo formulario (grupos de más de ICEBREAKER_MAX_MURO) ----
    url = formulario_url or formulario_presentacion_url(course_key)
    pendiente = not url.lower().startswith("http")
    pide_txt = pide or "**expectativa del curso** o **tema de interés** (1 frase)"
    sub_txt = sub or "Formulario de Google · encuestas y Q&A en el propio Meet"
    items = consignas or [
        f"**Escanea el QR o abre:** {url}",
        f"**Qué respondes (1 min):** tu **nombre** + {pide_txt}. Entras con tu correo "
        "**@cun.edu.co**: no hay que crear cuenta ni instalar nada.",
        (f"**Ahora, ~3 min.** Somos {n}: no alcanzamos a presentarnos uno por uno, "
         "así que el formulario nos ordena."
         if n is not None else "**Ahora, ~3 min.** El formulario ordena las presentaciones."),
        "**Con eso, hoy:** el Docente lee en voz alta **5 o 6 respuestas** y proyecta el "
        "**resumen** del formulario — ahí quedan todas a la vista, no solo las leídas.",
        "**En vivo:** las preguntas y las votaciones van por la **encuesta** y el **Q&A** del "
        "propio Meet; no hay que abrir ninguna otra plataforma.",
        "**Después:** las respuestas quedan en una hoja que el Docente usa todo el periodo "
        "(equipos, ejemplos y seguimiento).",
    ]
    # Sin enlace real no hay nada que codificar en un QR: se muestra el pendiente.
    qr = qr_path or ("" if pendiente else _qr_formulario(course_key))
    return _icebreaker_render(
        prs, idx, sub_txt, items, url, qr, size=13,
        qr_pendiente=("Se genera cuando el formulario tenga enlace real."
                      if pendiente else None),
    )

def link_callout_slide(prs, title, headline, url, notes=None, idx=None):
    """Slide para remarcar UN enlace obligatorio: caja grande verde + URL en tipografía grande."""
    s = blank(prs); bg_white(s)
    top = title_block(s, title)
    y = top + 0.25
    rounded(s, MARGIN, y, CONTENT_W, 2.4, INFO)
    rect(s, MARGIN, y, 0.18, 2.4, VERDE)
    tf = textbox(s, MARGIN + 0.45, y + 0.25, CONTENT_W - 0.7, 0.7)
    _rich(tf.paragraphs[0], headline, 20, NAVY, bold=True)
    uf = textbox(s, MARGIN + 0.45, y + 1.05, CONTENT_W - 0.7, 0.9, anchor=MSO_ANCHOR.MIDDLE)
    p = uf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p.add_run(), url, 22, VERDE, bold=True)
    y2 = y + 2.65
    for i, note in enumerate(notes or []):
        p = textbox(s, MARGIN, y2 + i * 0.55, CONTENT_W, 0.5).paragraphs[0]
        _rich(p, note, 14, GRAY)
    footer_num(s, idx)
    return s

# ---------- CRONOGRAMA EN TARJETAS INICIO / CIERRE (preferido: más legible) ----------
def fechas_inicio_fin_slide(prs, title, blocks, holiday_marks=None, note=None, sub=None, idx=None):
    """Cronograma en tarjetas grandes: cada bloque muestra INICIO y CIERRE con fecha completa.
      - blocks: lista de dicts {label, start(date), end(date), pct(str opc.)}
      - holiday_marks: lista de (date, label) — se listan debajo como 'sin clase'
      - note: texto inferior (p. ej. link de registro de tutorías)
    """
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    n = max(1, len(blocks))
    gap = 0.22
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = 3.15
    y = top + 0.15
    palette = [NAVY, SLATE, RGBColor(0x4A, 0x7F, 0xB5), RGBColor(0x2E, 0x7D, 0x5B)]
    for i, b in enumerate(blocks):
        x = MARGIN + i * (card_w + gap)
        color = b.get("color") or palette[i % len(palette)]
        # cabecera de color
        rect(s, x, y, card_w, 0.72, color)
        tf = textbox(s, x + 0.08, y + 0.08, card_w - 0.16, 0.56, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        label = b["label"] + (f"  ·  {b['pct']}" if b.get("pct") else "")
        _run(p.add_run(), label, 14, WHITE, bold=True)
        # cuerpo blanco
        rect(s, x, y + 0.72, card_w, card_h - 0.72, ALT)
        # INICIO
        tf = textbox(s, x + 0.12, y + 0.95, card_w - 0.24, 0.9)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), "INICIO", 11, SOFT, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
        _run(p2.add_run(), b["start"].strftime("%d/%m/%Y"), 20, VERDE, bold=True)
        # separador
        rect(s, x + 0.35, y + 1.95, card_w - 0.7, 0.03, RGBColor(0xD0, 0xD0, 0xD0))
        # CIERRE
        tf = textbox(s, x + 0.12, y + 2.15, card_w - 0.24, 0.9)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), "CIERRE", 11, SOFT, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
        _run(p2.add_run(), b["end"].strftime("%d/%m/%Y"), 20, NAVY, bold=True)

    # festivos
    y2 = y + card_h + 0.18
    if holiday_marks:
        fest = "  ·  ".join(d.strftime("%d/%m") + " " + lab.replace("\n", " ") for d, lab in holiday_marks)
        tf = textbox(s, MARGIN, y2, CONTENT_W, 0.35)
        _rich(tf.paragraphs[0], "🚫 **Sin clase (festivos):** " + fest, 12, RED)
        y2 += 0.38
    if note:
        rounded(s, MARGIN, y2, CONTENT_W, 0.7, ACLAR)
        tf = textbox(s, MARGIN + 0.2, y2 + 0.12, CONTENT_W - 0.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
        _rich(tf.paragraphs[0], note, 12, GRAY)
    footer_num(s, idx)
    return s

# ---------- LÍNEA DE TIEMPO (legacy — preferir fechas_inicio_fin_slide) ----------
def timeline_slide(prs, title, start_date, end_date, blocks, holiday_marks=None,
                    session_dots=None, sub=None, idx=None):
    """Línea de tiempo horizontal, con INICIO y CIERRE del curso resaltados en grande en los
    extremos, bloques con su rango de fechas bien separado y legible, y festivos con una guía
    vertical propia (fila separada, sin cruzarse con las fechas de los bloques).
      - blocks: lista de dicts {label, start(date), end(date), pct(str opc.), color(RGBColor opc.)}
      - holiday_marks: lista de (date, label) -> franja roja + línea guía + etiqueta.
      - session_dots: lista de date -> punto verde sobre la barra (sesión de clase real).
    """
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    bar_y = top + 1.15
    bar_h = 0.5
    x0, x1 = MARGIN + 1.35, SW - MARGIN - 1.35   # deja espacio a los lados para las banderas INICIO/CIERRE
    total_days = (end_date - start_date).days or 1

    def xpos(d):
        frac = max(0.0, min(1.0, (d - start_date).days / total_days))
        return x0 + frac * (x1 - x0)

    # eje base
    rect(s, x0, bar_y, x1 - x0, bar_h, ALT)

    # ---- banderas grandes de INICIO / CIERRE en los extremos ----
    for d, label, align, tx in [
        (start_date, "INICIO", PP_ALIGN.RIGHT, x0 - 1.3),
        (end_date, "CIERRE", PP_ALIGN.LEFT, x1 + 0.08),
    ]:
        tf = textbox(s, tx, bar_y - 0.08, 1.25, bar_h + 0.16, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = align
        _run(p.add_run(), label, 12, NAVY, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = align
        _run(p2.add_run(), d.strftime("%d/%m/%Y"), 12, VERDE, bold=True)
        # marca vertical gruesa en el extremo de la barra
        mx = x0 if label == "INICIO" else x1
        rect(s, mx - 0.02, bar_y - 0.12, 0.04, bar_h + 0.24, NAVY)

    # ---- bloques (ACA1/ACA2/ACA3/coev/autoev…) ----
    palette = [NAVY, SLATE, RGBColor(0x4A, 0x7F, 0xB5), RGBColor(0x7A, 0xA6, 0xCF), SOFT]
    for i, b in enumerate(blocks):
        bx0, bx1 = xpos(b["start"]), xpos(b["end"])
        bw = max(bx1 - bx0, 0.05)
        color = b.get("color") or palette[i % len(palette)]
        rect(s, bx0, bar_y, bw, bar_h, color)
        # etiqueta (nombre + %) ENCIMA del bloque
        label = b["label"] + (f" ({b['pct']})" if b.get("pct") else "")
        tf = textbox(s, bx0 - 0.3, bar_y - 0.42, max(bw, 1.6) + 0.6, 0.36)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _rich(tf.paragraphs[0], label, 12.5, NAVY, bold=True)
        # fechas de inicio-fin del bloque, GRANDES y bien separadas, DEBAJO
        df = textbox(s, bx0 - 0.35, bar_y + bar_h + 0.12, max(bw, 1.7) + 0.7, 0.32)
        p = df.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _rich(p, b["start"].strftime("%d/%m") + "  →  " + b["end"].strftime("%d/%m"), 12, VERDE, bold=True)

    # ---- festivos: fila SEPARADA (más abajo) con línea guía propia, sin chocar con las fechas ----
    fy_line_top = bar_y + bar_h + 0.55   # inicio de la línea guía (deja aire tras la fila de fechas)
    fy_label = fy_line_top + 0.28
    for d, label in (holiday_marks or []):
        fx = xpos(d)
        # línea guía vertical delgada desde la barra hasta la etiqueta
        rect(s, fx - 0.012, bar_y + bar_h, 0.024, fy_line_top - (bar_y + bar_h), RED)
        # marca corta sobre la barra
        rect(s, fx - 0.02, bar_y - 0.04, 0.04, bar_h + 0.08, RED)
        tf = textbox(s, fx - 0.6, fy_label, 1.3, 0.55)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _rich(p, "🚫 " + d.strftime("%d/%m"), 9.5, RED, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        _rich(p2, label, 8.5, RED)

    # ---- sesiones de clase reales: punto verde sobre la barra ----
    for item in (session_dots or []):
        d, label = item if isinstance(item, tuple) else (item, "")
        dx = xpos(d)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx - 0.055), Inches(bar_y + bar_h / 2 - 0.055),
                                  Inches(0.11), Inches(0.11))
        dot.fill.solid(); dot.fill.fore_color.rgb = VERDE; dot.line.color.rgb = WHITE
        dot.line.width = Pt(1); dot.shadow.inherit = False

    # ---- leyenda ----
    ly = SH - 0.5
    lx = MARGIN
    for color, text in [(VERDE, "Sesión de clase"), (RED, "Festivo (sin clase)"), (NAVY, "Inicio / cierre del bloque")]:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(ly), Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background(); dot.shadow.inherit = False
        _rich(textbox(s, lx + 0.20, ly - 0.03, 2.6, 0.25).paragraphs[0], text, 10, GRAY)
        lx += 2.7

    footer_num(s, idx)
    return s

# ---------- Portadas ----------
def course_cover(prs, materia, subtitulo, meta_lines):
    """Portada del curso: banda navy con el nombre de la asignatura + metadatos + logo real.

    El logo va arriba-izquierda y el título SIEMPRE arranca por debajo de él
    (no se superponen, aunque el título ocupe dos líneas).
    """
    BAND_H = 3.2
    s = blank(prs); bg_white(s)
    rect(s, 0, 0, SW, BAND_H, NAVY)
    rect(s, 0, BAND_H, SW, 0.08, VERDE)   # filete verde de marca bajo la banda
    # 1) Logo primero: fija el techo del área de título.
    logo_bottom = add_logo(s, width=1.45, corner="left-top", mt=0.3, mr=0.5)
    # 2) Título en la franja libre bajo el logo, centrado verticalmente.
    t_top = logo_bottom + 0.12
    t_h = max(0.8, BAND_H - t_top - 0.16)
    tf = textbox(s, MARGIN, t_top, CONTENT_W, t_h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    # Reserva ~0.44" para el subtítulo; el resto define cuántas líneas caben.
    avail = t_h - (0.44 if subtitulo else 0.0)
    max_lines = max(1, int(avail / 0.58))
    _run(p.add_run(), materia, _fit_title_size(materia, CONTENT_W, max_lines=max_lines),
         WHITE, bold=True)
    if subtitulo:
        ps = tf.add_paragraph(); ps.alignment = PP_ALIGN.CENTER; ps.space_before = Pt(8)
        _run(ps.add_run(), subtitulo, 16, RGBColor(0xCF, 0xDD, 0xEE))
    tm = textbox(s, MARGIN, 3.65, CONTENT_W, 2.9)
    for i, ln in enumerate(meta_lines or []):
        p = tm.paragraphs[0] if i == 0 else tm.add_paragraph()
        p.space_after = Pt(8); _rich(p, ln, 15, GRAY)
    footer_num(s, None)
    return s

def closing_slide(prs, big, lines, accent=None):
    s = blank(prs); bg_white(s)
    rect(s, 0, 0, SW, 0.18, VERDE)
    add_logo(s, width=1.5, corner="right-bottom", mt=0.3, mr=0.5)
    tf = textbox(s, MARGIN, 1.9, CONTENT_W, 1.2, anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _rich(tf.paragraphs[0], big, 34, NAVY, bold=True)
    tl = textbox(s, MARGIN, 3.3, CONTENT_W, 1.8)
    for i, ln in enumerate(lines):
        p = tl.paragraphs[0] if i == 0 else tl.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6); _rich(p, ln, 16, GRAY)
    if accent:
        ta = textbox(s, MARGIN, 5.2, CONTENT_W, 0.7); ta.paragraphs[0].alignment = PP_ALIGN.CENTER
        _rich(ta.paragraphs[0], accent, 20, VERDE, bold=True)
    footer_num(s, None)
    return s

def new_prs():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN)); prs.slide_height = Emu(int(SH * EMU_IN))
    cp = prs.core_properties
    cp.author = "CUN"; cp.last_modified_by = "CUN"
    cp.title = ""; cp.comments = ""
    return prs
