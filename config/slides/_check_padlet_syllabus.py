# -*- coding: utf-8 -*-
"""Auditoría de las decks YA GENERADAS: rompehielos correcto + terminología.

Qué comprueba, deck por deck (`Presentacion del Curso*.pptx` y `Sesion 01*/Presentacion.pptx`):

1. **Rompehielos según el TAMAÑO del grupo.** Desde el 2026-08-11 el Padlet dejó de ser
   «el mismo tablero en los 5 cursos»: solo lo usan los grupos que todavía se leen
   enteros (≤ `ICEBREAKER_MAX_MURO` estudiantes) y los demás van con el formulario de
   Google + encuestas/Q&A de Meet. Este chequeo NO trae una lista de cursos escrita a
   mano: le pregunta el modo a `cun_slides_engine.modo_rompehielos()`, que lo deriva de
   la matrícula real (roster de CDigital). Así, el día que cambie una matrícula, cambia
   la regla y cambia lo que aquí se exige.
   - Deck con slide de rompehielos (la Presentación del Curso): tiene que traer el del
     modo que le toca y **ninguna huella del otro**.
   - Deck de Sesión 01: no repite el rompehielos, pero tampoco puede mandar al muro a un
     curso que ya no lo usa (era el caso al hacer este cambio: las 5 decks de Sesión 01
     traían el QR del Padlet).
2. **Terminología:** «Syllabus», nunca «sílabo» (ver memoria del proyecto).
3. **Restos de versiones anteriores:** «Clear posts», «3 padlets», IdeaBoardz.

Solo audita: no regenera nada. Cada CHK dice qué build hay que volver a correr.
Uso: `python config/slides/_check_padlet_syllabus.py` (código de salida 1 si hay CHK).
"""
import os
import re
import sys
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from cun_slides_engine import (  # noqa: E402
    ICEBREAKER_MAX_MURO,
    MODO_FORMULARIO,
    MODO_MURO,
    PADLET_PRESENTACION_URL,
    contar_estudiantes,
    modo_rompehielos,
)

sys.stdout.reconfigure(encoding="utf-8")
# Raíz derivada del propio archivo (antes «G:\Mi unidad\…», que rompe cuando Google
# Drive monta la unidad en inglés «G:\My Drive»). Corregido 2026-08-09.
root = Path(__file__).resolve().parents[2]
CURSOS = ("proyecto1", "creatividad", "investigacion", "tg2", "tg3")

# Marca del rompehielos de grupo grande en el texto de la slide. Sirve tanto con el
# enlace real como con el marcador de posición: en los dos casos la slide habla de un
# «formulario» (y, mientras esté pendiente, del «[URL Formulario Preséntate …]»).
RX_FORMULARIO = re.compile(r"formulario", re.IGNORECASE)
RX_ROMPEHIELOS = re.compile(r"PRES[EÉ]NTATE|ROMPEHIELOS", re.IGNORECASE)


def _carpetas_de_curso() -> dict[str, Path]:
    """`{clave de curso: carpeta de la asignatura}` — para saber de quién es cada deck."""
    try:
        sys.path.insert(0, str(root / "config" / "cursos"))
        from carga_academica import course_dir  # noqa: E402
    except Exception:
        return {}
    out = {}
    for key in CURSOS:
        try:
            out[key] = Path(course_dir(key)).resolve()
        except Exception:
            continue
    return out


CARPETAS = _carpetas_de_curso()


def curso_de(path: Path) -> str | None:
    """Clave del curso al que pertenece la deck (por la carpeta de la asignatura)."""
    p = path.resolve()
    for key, carpeta in CARPETAS.items():
        try:
            p.relative_to(carpeta)
        except ValueError:
            continue
        return key
    return None


def textos_de(path: Path) -> tuple[str, list[str]]:
    """(texto completo de la deck, textos de las slides de rompehielos)."""
    prs = Presentation(str(path))
    todo, rompehielos = [], []
    for slide in prs.slides:
        blob = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        todo.append(blob)
        if RX_ROMPEHIELOS.search(blob):
            rompehielos.append(blob)
    return "\n".join(todo), rompehielos


def auditar(path: Path, key: str | None = None) -> tuple[str | None, int | None, list[str]]:
    """(modo esperado, matrícula, problemas) de una deck. `key` fuerza el curso (pruebas)."""
    text, slides_rh = textos_de(path)
    key = key or curso_de(path)
    n = contar_estudiantes(key) if key else None
    esperado = modo_rompehielos(key) if key else None
    n_txt = n if n is not None else "?"

    problemas: list[str] = []
    has_pad = PADLET_PRESENTACION_URL in text
    has_form = any(RX_FORMULARIO.search(b) for b in slides_rh)

    if key is None:
        problemas.append("deck fuera de las carpetas de los 5 cursos: no sé qué rompehielos le toca")
    elif slides_rh:
        # Es la deck que hace el rompehielos (la Presentación del Curso).
        if esperado == MODO_MURO:
            if not has_pad:
                problemas.append(f"grupo de {n_txt} (≤ {ICEBREAKER_MAX_MURO}): falta el muro de Padlet")
            if has_form:
                problemas.append("mezcla muro y formulario en la misma deck")
        else:
            if has_pad:
                problemas.append(f"grupo de {n_txt}: sigue mandando al Padlet")
            if not has_form:
                problemas.append(
                    f"grupo de {n_txt} (> {ICEBREAKER_MAX_MURO}): falta el rompehielos por formulario"
                )
    elif esperado == MODO_FORMULARIO and has_pad:
        # Deck sin slide de rompehielos que aun así deja el enlace del muro a la vista.
        problemas.append(f"grupo de {n_txt}: resto del Padlet en una deck sin rompehielos")

    sil = re.findall(r"[Ss][íi]labo|SÍLABO", text, flags=re.IGNORECASE)
    if sil:
        problemas.append(f"terminología: {sil[:3]!r} (se dice «Syllabus»)")
    if "Clear posts" in text or "3 padlets" in text:
        problemas.append("resto de versión anterior: «Clear posts» / «3 padlets»")
    if "IdeaBoardz" in text:
        problemas.append("resto de versión anterior: IdeaBoardz")
    return esperado, n, problemas


def main() -> int:
    paths = sorted(
        set(list(root.rglob("Presentacion del Curso*.pptx"))
            + list(root.rglob("Sesion 01*/Presentacion.pptx"))),
        key=lambda p: str(p),
    )
    fallos = revisadas = 0
    for p in paths:
        # "Temas" = carpetas legacy de PPTX sueltos; "_Archivo obsoleto" = material ya
        # retirado (no debe generar falsas alarmas de esta auditoría).
        if "Temas" in str(p) or "_Archivo obsoleto" in str(p):
            continue
        revisadas += 1
        esperado, n, problemas = auditar(p)
        fallos += bool(problemas)
        modo_txt = f"{esperado or '—'}" + (f"/{n}" if n is not None else "")
        print(f"{'OK ' if not problemas else 'CHK'} modo={modo_txt:<14} | {p.relative_to(root)}")
        for prob in problemas:
            print(f"      → {prob}")

    print(
        f"\n{revisadas} decks revisadas · {fallos} con hallazgos. "
        "Regenerar: `python config/slides/build_all_course_presentations.py` "
        "(Presentación del Curso) · `python config/slides/build_sesion_material.py` (Sesión NN)."
    )
    return 1 if fallos else 0


if __name__ == "__main__":
    sys.exit(main())
