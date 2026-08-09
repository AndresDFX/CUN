# -*- coding: utf-8 -*-
"""
Motor de diapositivas .pptx para el curso "Herramientas de Google Cloud Platform"
(UMNG para la UGPP). Replica y MEJORA la identidad de los decks Nivel 1-4:
  - Portadas: banda DORADA (#B08840, color UMNG) con título blanco centrado + barra
    tricolor de Colombia abajo.
  - Contenido: fondo blanco, logos UGPP + UMNG arriba-derecha, título dorado con regla,
    viñetas, bloques de código SQL, tablas de "resultado esperado".
  - Extras para este proyecto: slides mixtas texto+pantallazo, y slides de FAQ (docente).

Uso:
    from gcp_slides_engine import *
    prs = new_prs()
    cover(prs, "SESIÓN 8", "Refuerzo de SQL en BigQuery", "Guía del estudiante · práctica")
    content_slide(prs, "Objetivo de hoy", ["...", ("sub", 1)])
    step_slide(prs, "Paso 1 · Crear tu proyecto", ["..."], "ruta/pantallazo.png")
    code_slide(prs, "Consulta: ventas por ciudad", ["SELECT ciudad,", "  SUM(valor) AS total", "FROM ..."], nota="...")
    table_content(prs, "Resultado esperado", ["Ciudad","Total"], [["Bogotá","5100"]])
    faq_slide(prs, [("¿Pregunta?","Respuesta...")])
    prs.save("salida.pptx")
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Paleta (UMNG/UGPP) ----------
GOLD   = RGBColor(0xB0, 0x88, 0x40)   # dorado UMNG (banners/acentos)
GOLD_D = RGBColor(0x8A, 0x69, 0x2E)   # dorado oscuro (títulos)
DARK   = RGBColor(0x22, 0x22, 0x22)   # texto casi negro
GRAY   = RGBColor(0x4A, 0x4A, 0x49)
SOFT   = RGBColor(0x8A, 0x8A, 0x8A)
BANNER = RGBColor(0xF3, 0xEE, 0xE2)   # crema muy claro (cajas suaves)
LINK   = RGBColor(0x1A, 0x73, 0xE8)   # azul Google (enlaces)
CODEBG = RGBColor(0x1E, 0x1E, 0x2E)   # fondo de código (oscuro)
CODEFG = RGBColor(0xE6, 0xE6, 0xE6)
CODEKW = RGBColor(0xF0, 0xB4, 0x5A)   # palabra clave SQL (dorado claro sobre oscuro)
ALT    = RGBColor(0xEC, 0xE6, 0xD7)   # fila alterna de tabla
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
# Bandera de Colombia
FLAG_Y = RGBColor(0xFC, 0xD1, 0x16); FLAG_B = RGBColor(0x00, 0x33, 0x8D); FLAG_R = RGBColor(0xCE, 0x11, 0x26)
FONT = "Barlow"; MONO = "Consolas"

EMU_IN = 914400
SW, SH = 13.333, 7.5
MARGIN = 0.75
CONTENT_W = SW - 2 * MARGIN
_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "gcp")
UGPP = os.path.join(_ASSETS, "ugpp_logo.jpeg")
UMNG = os.path.join(_ASSETS, "umng_logo.png")
_FOOTER = ""

def set_footer(t):
    global _FOOTER; _FOOTER = t

# ---------- Texto enriquecido: **negrita** · @@dorado@@ · `código` ----------
def _run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text; f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font; f.color.rgb = color

def _rich(paragraph, markup, size, base_color, bold=False):
    for part in re.split(r'(@@.*?@@|\*\*.*?\*\*|\*[^*\n]+?\*|`[^`]+`)', markup):
        if not part:
            continue
        r = paragraph.add_run()
        if part.startswith("@@") and part.endswith("@@"):
            _run(r, part[2:-2], size, GOLD_D, bold=True)
        elif part.startswith("**") and part.endswith("**"):
            _run(r, part[2:-2], size, base_color, bold=True)
        elif len(part) > 2 and part.startswith("*") and part.endswith("*"):
            _run(r, part[1:-1], size, base_color, bold=bold, italic=True)
        elif part.startswith("`") and part.endswith("`"):
            _run(r, part[1:-1], size - 1, DARK, font=MONO)
        else:
            _run(r, part, size, base_color, bold=bold)

# ---------- Primitivas ----------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def _img_size(path, default=(4, 3)):
    try:
        from PIL import Image as _I; return _I.open(path).size
    except Exception:
        return default

def flag_bar(slide, cx, y, w=1.5, h=0.16):
    """Barra tricolor de Colombia centrada en cx."""
    x = cx - w / 2
    rect(slide, x, y, w * 0.5, h, FLAG_Y)
    rect(slide, x + w * 0.5, y, w * 0.25, h, FLAG_B)
    rect(slide, x + w * 0.75, y, w * 0.25, h, FLAG_R)

def dual_logos(slide, y=0.28, h=0.62):
    """UGPP + separador + UMNG arriba-derecha (co-marca)."""
    x = SW - MARGIN
    if os.path.exists(UMNG):
        iw, ih = _img_size(UMNG); w = h * iw / ih
        x -= w; slide.shapes.add_picture(UMNG, Inches(x), Inches(y), height=Inches(h))
        x -= 0.18
    rect(slide, x, y + 0.05, 0.012, h - 0.1, SOFT); x -= 0.18
    if os.path.exists(UGPP):
        iw, ih = _img_size(UGPP); w = h * iw / ih
        x -= w; slide.shapes.add_picture(UGPP, Inches(x), Inches(y), height=Inches(h))

def footer_num(slide, idx):
    if idx is not None:
        tf = textbox(slide, MARGIN, SH - 0.42, 4.0, 0.3)
        _run(tf.paragraphs[0].add_run(), (_FOOTER + ("  ·  " if _FOOTER else "") + str(idx)), 9, SOFT)

def _title(slide, title, sub=None):
    """Encabezado de contenido: logos + título dorado + regla."""
    dual_logos(slide)
    tf = textbox(slide, MARGIN, 1.15, CONTENT_W - 2.6, 0.8)
    _rich(tf.paragraphs[0], title, 27, GOLD_D, bold=True)
    rect(slide, MARGIN, 1.95, CONTENT_W, 0.028, GOLD)
    top = 2.2
    if sub:
        ts = textbox(slide, MARGIN, top, CONTENT_W, 0.4)
        _rich(ts.paragraphs[0], sub, 14, GRAY); top += 0.5
    return top

def bullets(slide, items, top, size=16, left=None, width=None):
    left = MARGIN if left is None else left
    width = CONTENT_W if width is None else width
    tf = textbox(slide, left, top, width, SH - top - 0.55)
    for i, it in enumerate(items):
        text, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(9 if lvl == 0 else 5)
        _rich(p, ("●   " if lvl == 0 else "–   ") + text, size - (1 if lvl else 0), GRAY)
    return tf

# ---------- Slides ----------
def cover(prs, kicker, titulo, subtitulo=None, meta=None):
    s = blank(prs); bg(s, WHITE)
    rect(s, 0, 1.0, SW, 5.1, GOLD)   # banda dorada
    # kicker + título + subtítulo, centrados en la banda
    tk = textbox(s, MARGIN, 2.15, CONTENT_W, 0.6); tk.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tk.paragraphs[0].add_run(), kicker, 20, WHITE, bold=True)
    tt = textbox(s, MARGIN, 2.9, CONTENT_W, 1.5, anchor=MSO_ANCHOR.TOP); tt.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tt.paragraphs[0].add_run(), titulo, 36, WHITE, bold=True)
    if subtitulo:
        ps = tt.add_paragraph(); ps.alignment = PP_ALIGN.CENTER; ps.space_before = Pt(8)
        _run(ps.add_run(), subtitulo, 20, RGBColor(0xF3, 0xEE, 0xE2))
    dual_logos(s, y=0.24, h=0.6)
    flag_bar(s, SW / 2, 6.45, w=1.7, h=0.16)
    # (Se quitó la línea de texto inferior "UGPP · Universidad Militar Nueva Granada" — no va en la portada.)
    return s

def content_slide(prs, title, items, sub=None, size=16, idx=None):
    s = blank(prs); bg(s)
    top = _title(s, title, sub)
    bullets(s, items, top + 0.05, size=size)
    footer_num(s, idx); return s

def step_slide(prs, title, items, image_path, sub=None, size=15, idx=None):
    """Paso práctico: instrucciones a la izquierda + pantallazo (enmarcado) a la derecha."""
    s = blank(prs); bg(s)
    top = _title(s, title, sub)
    text_w = 6.2
    bullets(s, items, top + 0.05, size=size, left=MARGIN, width=text_w)
    if image_path and os.path.exists(image_path):
        px = MARGIN + text_w + 0.35; pw = SW - MARGIN - px
        py = top + 0.05; ph = SH - py - 0.55
        iw, ih = _img_size(image_path); ratio = iw / ih
        w = pw; h = w / ratio
        if h > ph:
            h = ph; w = h * ratio
        fx = px + (pw - w) / 2; fy = py + (ph - h) / 2
        rect(s, fx - 0.04, fy - 0.04, w + 0.08, h + 0.08, BANNER, line=GOLD)  # marco
        s.shapes.add_picture(image_path, Inches(fx), Inches(fy), Inches(w), Inches(h))
    else:
        px = MARGIN + text_w + 0.35; pw = SW - MARGIN - px
        ph = SH - (top + 0.05) - 0.55
        r = rect(s, px, top + 0.05, pw, ph, BANNER, line=GOLD)
        t = textbox(s, px, top + ph / 2 - 0.3, pw, 0.6); t.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(t.paragraphs[0].add_run(), "📸 [pantallazo]", 12, SOFT, italic=True)
    footer_num(s, idx); return s

def code_slide(prs, title, code_lines, nota=None, sub=None, idx=None):
    """Bloque de código SQL (caja oscura, palabras clave resaltadas)."""
    s = blank(prs); bg(s)
    top = _title(s, title, sub)
    kw = re.compile(r'\b(SELECT|FROM|WHERE|GROUP BY|ORDER BY|JOIN|LEFT|INNER|ON|AS|SUM|AVG|COUNT|MIN|MAX|HAVING|LIMIT|DESC|ASC|AND|OR|NOT|WITH|CREATE|MODEL|OVER|PARTITION BY|CASE|WHEN|THEN|END|DISTINCT)\b')
    n = len(code_lines)
    bh = min(SH - top - 1.0, 0.34 * n + 0.4)
    box = rect(s, MARGIN, top + 0.05, CONTENT_W, bh, CODEBG)
    tf = textbox(s, MARGIN + 0.25, top + 0.2, CONTENT_W - 0.5, bh - 0.3)
    for i, ln in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        pos = 0
        for m in kw.finditer(ln):
            if m.start() > pos:
                _run(p.add_run(), ln[pos:m.start()], 13, CODEFG, font=MONO)
            _run(p.add_run(), m.group(0), 13, CODEKW, bold=True, font=MONO)
            pos = m.end()
        if pos < len(ln):
            _run(p.add_run(), ln[pos:], 13, CODEFG, font=MONO)
        if not ln:
            _run(p.add_run(), " ", 13, CODEFG, font=MONO)
    if nota:
        _rich(textbox(s, MARGIN, top + 0.05 + bh + 0.15, CONTENT_W, 0.7).paragraphs[0], nota, 13, GRAY)
    footer_num(s, idx); return s

def table_content(prs, title, headers, rows, nota=None, col_w=None, sub=None, idx=None):
    s = blank(prs); bg(s)
    top = _title(s, title, sub)
    ncol = len(headers); nrow = len(rows) + 1
    col_w = col_w or [CONTENT_W / ncol] * ncol
    tot = sum(col_w); col_w = [w / tot * CONTENT_W for w in col_w]
    height = min(SH - top - 0.9, 0.5 * nrow)
    table = s.shapes.add_table(nrow, ncol, Inches(MARGIN), Inches(top + 0.05), Inches(CONTENT_W), Inches(height)).table
    tblPr = table._tbl.tblPr
    for c in list(tblPr):
        if c.tag == qn('a:tableStyleId'):
            tblPr.remove(c)
    sid = tblPr.makeelement(qn('a:tableStyleId'), {}); sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'; tblPr.append(sid)
    table.first_row = False; table.horz_banding = False
    for j, w in enumerate(col_w):
        table.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = table.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = GOLD
        c.vertical_anchor = MSO_ANCHOR.MIDDLE; c.margin_left = Inches(0.1); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
        _run(c.text_frame.paragraphs[0].add_run(), h, 12, WHITE, bold=True)
    for i, row in enumerate(rows):
        rc = WHITE if i % 2 == 0 else ALT
        for j in range(ncol):
            c = table.cell(i + 1, j); c.fill.solid(); c.fill.fore_color.rgb = rc
            c.vertical_anchor = MSO_ANCHOR.MIDDLE; c.margin_left = Inches(0.1); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            _rich(c.text_frame.paragraphs[0], row[j] if j < len(row) else "", 11, GRAY)
    if nota:
        _rich(textbox(s, MARGIN, top + 0.05 + height + 0.15, CONTENT_W, 0.7).paragraphs[0], nota, 12, GRAY)
    footer_num(s, idx); return s

def faq_slide(prs, faqs, title="Preguntas frecuentes de la práctica", image_path=None, sub=None, idx=None):
    """Slide de FAQ (docente): pregunta en dorado + respuesta. Opcional pantallazo a la derecha."""
    s = blank(prs); bg(s)
    top = _title(s, title, sub)
    width = CONTENT_W if not image_path else 6.6
    tf = textbox(s, MARGIN, top + 0.05, width, SH - top - 0.6)
    first = True
    for q, a in faqs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(0 if p is tf.paragraphs[0] else 10); p.space_after = Pt(2)
        _rich(p, "❓ " + q, 15, GOLD_D, bold=True)
        pa = tf.add_paragraph(); pa.space_after = Pt(4)
        _rich(pa, a, 13, GRAY)
    if image_path and os.path.exists(image_path):
        px = MARGIN + width + 0.3; pw = SW - MARGIN - px; py = top + 0.05; ph = SH - py - 0.6
        iw, ih = _img_size(image_path); ratio = iw / ih; w = pw; h = w / ratio
        if h > ph: h = ph; w = h * ratio
        fx = px + (pw - w) / 2; fy = py + (ph - h) / 2
        rect(s, fx - 0.04, fy - 0.04, w + 0.08, h + 0.08, BANNER, line=GOLD)
        s.shapes.add_picture(image_path, Inches(fx), Inches(fy), Inches(w), Inches(h))
    footer_num(s, idx); return s

def closing(prs, mensaje, lineas=None):
    s = blank(prs); bg(s, WHITE)
    rect(s, 0, 2.6, SW, 2.3, GOLD)
    tf = textbox(s, MARGIN, 3.0, CONTENT_W, 1.0, anchor=MSO_ANCHOR.MIDDLE); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tf.paragraphs[0].add_run(), mensaje, 30, WHITE, bold=True)
    if lineas:
        tl = textbox(s, MARGIN, 4.0, CONTENT_W, 0.8);
        for i, ln in enumerate(lineas):
            p = tl.paragraphs[0] if i == 0 else tl.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            _run(p.add_run(), ln, 14, RGBColor(0xF3, 0xEE, 0xE2))
    dual_logos(s, y=0.24, h=0.6); flag_bar(s, SW / 2, 5.15, w=1.7, h=0.16)
    return s

def new_prs():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN)); prs.slide_height = Emu(int(SH * EMU_IN))
    cp = prs.core_properties; cp.author = "UGPP · UMNG"; cp.title = ""
    return prs
