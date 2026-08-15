# -*- coding: utf-8 -*-
"""Extrae el texto de un PDF, DOCX o PPTX con marcas de página/diapositiva.

Está pensado para leer trabajos de grado de 60-120 páginas sin gastar el
presupuesto de contexto en imágenes: devuelve texto plano con un separador
`===== [p. N] =====` por página, de modo que cualquier afirmación sobre el
documento pueda citarse con número de página verificable.

Uso:
    python config/evaluaciones/extraer_texto.py "ruta/al/documento.pdf"
    python config/evaluaciones/extraer_texto.py "doc.pdf" --paginas 1-20
    python config/evaluaciones/extraer_texto.py "doc.pdf" --indice
    python config/evaluaciones/extraer_texto.py "doc.pdf" --buscar "muestra,Turnitin"
"""
from __future__ import annotations

import argparse
import os
import re
import sys

sys.stdout.reconfigure(encoding="utf-8")


def _rango(txt: str, tope: int) -> range:
    if not txt:
        return range(1, tope + 1)
    m = re.match(r"^\s*(\d+)\s*(?:-\s*(\d+))?\s*$", txt)
    if not m:
        raise SystemExit(f"--paginas no entendido: {txt!r} (usa 5 o 5-20)")
    a = int(m.group(1))
    b = int(m.group(2) or m.group(1))
    return range(max(1, a), min(tope, b) + 1)


def paginas_pdf(ruta: str) -> list[str]:
    from pypdf import PdfReader

    r = PdfReader(ruta)
    fuera = []
    for p in r.pages:
        try:
            fuera.append(p.extract_text() or "")
        except Exception as e:  # una página corrupta no debe tumbar el informe
            fuera.append(f"[[página ilegible: {type(e).__name__}]]")
    return fuera


def paginas_docx(ruta: str) -> list[str]:
    from docx import Document

    d = Document(ruta)
    trozos = [p.text for p in d.paragraphs]
    for t in d.tables:
        for f in t.rows:
            celdas = [c.text.strip() for c in f.cells]
            if any(celdas):
                trozos.append(" | ".join(celdas))
    return ["\n".join(trozos)]


def paginas_pptx(ruta: str) -> list[str]:
    from pptx import Presentation

    pr = Presentation(ruta)
    fuera = []
    for s in pr.slides:
        t = []
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t.append(sh.text_frame.text)
            if getattr(sh, "has_table", False):
                for f in sh.table.rows:
                    t.append(" | ".join(c.text.strip() for c in f.cells))
        if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip():
            t.append("[NOTAS] " + s.notes_slide.notes_text_frame.text)
        fuera.append("\n".join(t))
    return fuera


def leer(ruta: str) -> tuple[list[str], str]:
    ext = os.path.splitext(ruta)[1].lower()
    if ext == ".pdf":
        return paginas_pdf(ruta), "p."
    if ext == ".docx":
        return paginas_docx(ruta), "doc"
    if ext == ".pptx":
        return paginas_pptx(ruta), "slide"
    raise SystemExit(f"extensión no soportada: {ext} (pdf, docx, pptx)")


def main() -> None:
    ap = argparse.ArgumentParser(description="Texto de un PDF/DOCX/PPTX con marcas de página.")
    ap.add_argument("ruta")
    ap.add_argument("--paginas", default="", help="5 o 5-20; por defecto todas")
    ap.add_argument("--indice", action="store_true", help="solo la primera línea útil de cada página")
    ap.add_argument("--buscar", default="", help="términos separados por coma; imprime página y línea")
    a = ap.parse_args()

    if not os.path.exists(a.ruta):
        raise SystemExit(f"no existe: {a.ruta}")

    pags, unidad = leer(a.ruta)
    print(f"### {os.path.basename(a.ruta)} — {len(pags)} {unidad if unidad != 'p.' else 'páginas'}")

    if a.buscar:
        terminos = [t.strip().lower() for t in a.buscar.split(",") if t.strip()]
        for i, txt in enumerate(pags, 1):
            for linea in txt.splitlines():
                bajo = linea.lower()
                if any(t in bajo for t in terminos):
                    print(f"[{unidad} {i}] {linea.strip()[:220]}")
        return

    # Encabezados y pies repetidos: una línea que sale en más de un tercio de las
    # páginas es membrete, no contenido, y con --indice tapa el título real.
    boiler: set[str] = set()
    if a.indice and len(pags) > 3:
        from collections import Counter
        c = Counter(l.strip() for txt in pags for l in txt.splitlines() if len(l.strip()) > 3)
        boiler = {l for l, n in c.items() if n > len(pags) / 3}

    for i in _rango(a.paginas, len(pags)):
        txt = pags[i - 1]
        if a.indice:
            utiles = [l.strip() for l in txt.splitlines()
                      if len(l.strip()) > 3 and l.strip() not in boiler]
            print(f"[{unidad} {i}] {' · '.join(utiles[:2])[:160] or '(sin contenido propio)'}")
        else:
            print(f"\n===== [{unidad} {i}] =====")
            print(txt.strip() or "(sin texto extraíble — puede ser imagen escaneada)")


if __name__ == "__main__":
    main()
