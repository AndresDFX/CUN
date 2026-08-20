# -*- coding: utf-8 -*-
"""Verifica que el mapa «🗺️ Slides de esta presentación» de cada guion diga los títulos
que el deck REALMENTE tiene.

Por qué existe: varios generadores (`_regen_guiones_pregrado.py`,
`_regen_guiones_proyecto1.py`) escriben esa tabla **a mano**, no leyéndola del `.pptx`.
Cuando el deck se reescribe, la tabla queda mintiendo y el Docente busca en pantalla una
slide con un título que ya no existe — el fallo silencioso más caro de este material,
porque solo se descubre en clase.

    python config/slides/verificar_mapas_slides.py          # todos los cursos
    python config/slides/verificar_mapas_slides.py --v      # muestra también los que sí cuadran

Sale con código 1 si hay algún desfase, para poder encadenarlo en un check.

Dos decisiones para no gritar en falso:
* Se compara contra **todos** los cuadros de texto de la slide, no solo contra
  `shapes.title`. En estos decks la slide del rompehielos tiene «slido.com» en letra
  gigante y `titulos_pptx()` la toma por título; el encabezado visible es otro.
* Una fila que abarca un rango (`| **8–9** |`) resume dos slides, así que basta con que
  cuadre el trozo anterior a los dos puntos.
"""
from __future__ import annotations

import glob
import os
import re
import sys
import unicodedata

sys.stdout.reconfigure(encoding="utf-8")
SLIDES = os.path.dirname(os.path.abspath(__file__))
RAIZ = os.path.abspath(os.path.join(SLIDES, "..", ".."))
sys.path.insert(0, SLIDES)

from pptx import Presentation  # noqa: E402

from guion_slides import deck_path  # noqa: E402

CURSOS = (
    "Pregrado/Creatividad y pensamiento innovador",
    "Pregrado/Investigacion en ciencia y tecnologia",
    "Pregrado/Trabajo de grado 2",
    "Pregrado/Trabajo de grado 3",
    "Especializacion/Proyecto I",
)

# Filas cuyo texto es DESCRIPTIVO a propósito (no pretende copiar el título del deck).
TOLERADOS = ("portada", "cierre —", "cierre -")

FILA = re.compile(r"^\|\s*\*\*(\d+)(?:\s*[–\-—]\s*(\d+))?\*\*\s*\|([^|]*)\|")


def candidatos(pptx_path: str) -> list[list[str]] | None:
    """Por cada slide, TODOS los primeros renglones de sus cuadros de texto."""
    if not pptx_path or not os.path.isfile(pptx_path):
        return None
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return None
    out = []
    for slide in prs.slides:
        textos = []
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            t = sh.text_frame.text.strip()
            if t:
                textos.append(t.splitlines()[0].strip())
        out.append(textos)
    return out


def limpio(s: str) -> str:
    """Compara sin acentos, sin negritas, sin la glosa entre paréntesis del final."""
    s = s.replace("**", "").strip()
    s = re.sub(r"\s*\([^()]*\)\s*$", "", s)          # «… (juego en Slido)» → «…»
    s = s.rstrip(" .…")
    s = unicodedata.normalize("NFKD", s)
    s = "".join(c for c in s if not unicodedata.combining(c))
    return re.sub(r"[^a-z0-9]+", " ", s.lower()).strip()


def casa(dicho: str, real: str, *, resumen: bool = False) -> bool:
    """Cuadra si uno es prefijo del otro: los generadores truncan títulos largos."""
    a, b = limpio(dicho), limpio(real)
    if not a or not b:
        return False
    if a == b or a.startswith(b) or b.startswith(a):
        return True
    if resumen:                                       # fila que resume un rango de slides
        a2, b2 = limpio(dicho.split(":")[0]), limpio(real.split(":")[0])
        return bool(a2) and bool(b2) and (a2 == b2 or a2.startswith(b2) or b2.startswith(a2))
    return False


def tabla(texto: str) -> list[str]:
    """Filas del mapa de slides — exige el encabezado «| Slide | Título en el PPTX |»."""
    lineas = texto.splitlines()
    for i, ln in enumerate(lineas):
        if "Slides de esta presentaci" not in ln:
            continue
        # ventana holgada: entre el título y la tabla puede haber una nota al margen
        for j in range(i + 1, min(i + 14, len(lineas))):
            cab = lineas[j].strip()
            if cab.startswith("|") and "Slide" in cab and "PPTX" in cab:
                filas = []
                for k in range(j + 2, len(lineas)):   # +2 salta el separador |:---:|
                    if not lineas[k].strip().startswith("|"):
                        break
                    filas.append(lineas[k])
                return filas
    return []


def inline(texto: str) -> list[tuple[int, int, str]]:
    """Referencias «**Slides:** 5 (TÍTULO) → 10 (OTRO)» de las fases → (ini, fin, título).

    Los paréntesis se cierran contando, no con regex: hay títulos que contienen
    «(cont.)» dentro del propio paréntesis.
    """
    out: list[tuple[int, int, str]] = []
    for ln in texto.splitlines():
        s = ln.strip()
        if not s.startswith("**Slides:**"):
            continue
        for trozo in s[len("**Slides:**"):].split("→"):
            trozo = trozo.strip()
            m = re.match(r"^\**(\d+)\**(?:\s*[–\-—]\s*\**(\d+)\**)?\s*\(", trozo)
            if not m:
                continue
            prof, ini_txt = 0, m.end() - 1
            for k in range(ini_txt, len(trozo)):
                prof += (trozo[k] == "(") - (trozo[k] == ")")
                if prof == 0:
                    out.append((int(m.group(1)), int(m.group(2) or m.group(1)),
                                trozo[ini_txt + 1:k].strip()))
                    break
    return out


def main() -> int:
    verboso = "--v" in sys.argv
    fallos = revisados = filas_ok = sin_tabla = 0
    ref_ok = ref_mal = 0

    for curso in CURSOS:
        carpeta = os.path.join(RAIZ, curso.replace("/", os.sep))
        for md in sorted(glob.glob(os.path.join(carpeta, "Docente", "Guiones", "Sesion *.md"))):
            label = os.path.splitext(os.path.basename(md))[0]
            reales = candidatos(deck_path(carpeta, label))
            if not reales:
                print(f"·  SIN DECK  {curso} · {label}")
                continue
            texto = open(md, encoding="utf-8").read()
            filas = tabla(texto)
            if not filas:
                sin_tabla += 1
            else:
                revisados += 1
            malas = []

            # --- referencias inline de las fases: «**Slides:** 5 (TÍTULO) → 10 (OTRO)»
            for ini, fin, dicho in inline(texto):
                if any(dicho.lower().startswith(t) for t in TOLERADOS):
                    continue
                pool = [t for i in range(ini, fin + 1) if 1 <= i <= len(reales)
                        for t in reales[i - 1]]
                if pool and any(casa(dicho, r, resumen=True) for r in pool):
                    ref_ok += 1
                else:
                    ref_mal += 1
                    encabezados = [reales[i - 1][0] for i in range(ini, fin + 1)
                                   if 1 <= i <= len(reales) and reales[i - 1]]
                    malas.append((f"{ini}*", dicho,
                                  " / ".join(encabezados) or "(el deck no tiene esa slide)"))

            for linea in filas:
                m = FILA.match(linea.strip())
                if not m:
                    continue
                dicho = m.group(3).strip()
                if any(dicho.lower().startswith(t) for t in TOLERADOS):
                    continue
                ini, fin = int(m.group(1)), int(m.group(2) or m.group(1))
                rango = fin > ini
                pool = [t for i in range(ini, fin + 1) if 1 <= i <= len(reales)
                        for t in reales[i - 1]]
                if not pool:
                    malas.append((str(ini), dicho, "(el deck no tiene esa slide)"))
                elif not any(casa(dicho, r, resumen=rango) for r in pool):
                    encabezados = [reales[i - 1][0] for i in range(ini, fin + 1)
                                   if 1 <= i <= len(reales) and reales[i - 1]]
                    malas.append((str(ini), dicho, " / ".join(encabezados)))
                else:
                    filas_ok += 1
            if malas:
                fallos += len(malas)
                print(f"\n✗  {curso} · {label}   ({len(reales)} slides en el deck)")
                for n, dicho, real in malas:
                    print(f"     slide {n:>4s}   guion: {dicho}")
                    print(f"                  deck : {real}")
            elif verboso:
                print(f"✓  {curso} · {label}")

    print(f"\n{'=' * 90}\n{revisados} guiones con mapa  ·  {sin_tabla} sin mapa\n"
          f"tabla   : {filas_ok} filas correctas\n"
          f"inline  : {ref_ok} referencias «**Slides:** N (…)» correctas\n"
          f"DESFASES: {fallos}   (las filas marcadas con * son referencias inline)")
    return 1 if fallos else 0


if __name__ == "__main__":
    raise SystemExit(main())
