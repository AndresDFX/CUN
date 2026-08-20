# -*- coding: utf-8 -*-
"""Comprueba el refactor de talleres sobre los .pptx ya construidos, no sobre el JSON.

Lo que ve el estudiante es la deck, así que la prueba se hace ahí: que cada sesión con taller
traiga sus **dos** slides, que ninguna sea una «(cont.)» sin título propio, y que no quede ni una
frase que lo mande a un espacio del aula que no existe.
"""
import glob
import os
import re
import sys

sys.stdout.reconfigure(encoding="utf-8")
_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _HERE)
sys.path.insert(0, os.path.join(os.path.dirname(_HERE), "cursos"))
from pptx import Presentation  # noqa: E402

import talleres as T  # noqa: E402
from sesiones_cun import CDIGITAL_PLACEHOLDER, COURSES  # noqa: E402

# La carpeta la declara `COURSES["folder"]`: copiarla aquí a mano ya me costó cinco sesiones
# invisibles («Investigacion **en** ciencia y tecnologia»).
CARPETA = {k: c["folder"] for k, c in COURSES.items() if c.get("folder")}
# Las frases que el refactor tenía que eliminar de la vista del estudiante.
FALSAS = [CDIGITAL_PLACEHOLDER, "espacio de esa sesión", "espacio de esta sesión",
          "espacio de la sesión"]

fallos: list[str] = []
con_taller = 0

for curso, sub in sorted(CARPETA.items()):
    base = os.path.join(sub, "Clases")
    for carp in sorted(glob.glob(os.path.join(base, "Sesion *"))):
        m = re.search(r"Sesion (\d+)", os.path.basename(carp))
        pptx = os.path.join(carp, "Presentacion.pptx")
        if not m or not os.path.isfile(pptx):
            continue
        n = int(m.group(1))
        titulos, texto = [], []
        for sl in Presentation(pptx).slides:
            t = ""
            for sh in sl.shapes:
                if sh.has_text_frame:
                    texto.append(sh.text_frame.text)
                    if not t:
                        t = sh.text_frame.text.strip().splitlines()[0] if sh.text_frame.text.strip() else ""
            titulos.append(t)
        todo = "\n".join(texto)
        ref = "%s s%02d" % (curso, n)

        de_taller = [t for t in titulos if t.upper().startswith("TALLER")]
        if T.tiene(curso, n):
            con_taller += 1
            if len(de_taller) != 2:
                fallos.append("%s: %d slides de taller (deberían ser 2): %s"
                              % (ref, len(de_taller), de_taller))
            if any("(cont.)" in t for t in de_taller):
                fallos.append("%s: hay una slide de taller «(cont.)» sin título propio." % ref)
            e = T.entrada(curso, n)
            if e["archivo"] not in todo:
                fallos.append("%s: la deck no nombra el archivo «%s»." % (ref, e["archivo"]))
        elif de_taller:
            fallos.append("%s: tiene slide de taller pero no hay entrada en talleres.py: %s"
                          % (ref, de_taller))

        for f in FALSAS:
            if f in todo:
                fallos.append("%s: la deck todavía dice «%s»." % (ref, f))

print("%d sesiones con taller revisadas en sus .pptx" % con_taller)
print("%d fallos" % len(fallos))
for f in fallos:
    print("   ✗ " + f)
sys.exit(1 if fallos else 0)
