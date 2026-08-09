# -*- coding: utf-8 -*-
"""
Motor reutilizable para generar diapositivas .pptx con la identidad de Nueva América (FESNA).
Diseño alineado a la PLANTILLA BASE institucional y a los decks reales de curso (Internetworking):
  - Título centrado sobre BANDA GRIS clara (#F0F1F2) en la parte superior.
  - Viñetas con guion (–) en nivel 0 y círculo (●) en nivel 1.
  - Portada del curso: FOTO institucional LAVADA (19% opacidad) al lado derecho ·
    "¡Bienvenidos estudiantes!" en AZUL institucional sobre la foto · materia debajo ·
    datos del programa arriba-izq · "Empezamos a las…" abajo-izq · LOGO abajo-der.
  - Objetivos: foto a sangre al lado derecho con VELO AZUL (#0A499C ~45%) + tarjeta
    blanca a la izquierda con las viñetas y etiqueta blanca sobre la foto.
  - Gestores: flyer institucional real como imagen centrada (image_slide).
  - Portada de sesión: etiqueta grande "SESIÓN N" + nivel + tema + foto lavada + logo.

PALETA OFICIAL (de la guía de marca):
  Naranja #FD531E · Azul #0A499C (logo, bienvenida de portada y velo de fotos) ·
  Gris principal #4A4A49 · negros/grises de soporte. Tipografía Barlow.
  Naranja para resaltar (markup @@texto@@). El azul NO se usa en texto corriente ni tablas.

USO:
    from fesna_slides_engine import *
    set_footer("Nombre del curso")
    prs = new_prs()
    course_cover(prs, materia, leccion, welcome, program_lines, start_note)
    objectives_slide(prs, "Objetivos", ["Objetivo 1", "Objetivo 2"], idx=4)
    session_cover(prs, "SESIÓN 1", nivel, titulo, subtitulo, welcome, meta_lines)
    content_slide(prs, "Título", ["@@resaltado@@", ("sub", 1)], idx=2)
    table_content(prs, "Título", ["A","B"], [["x","y"]], idx=3)
    image_slide(prs, "GESTORES", GESTORES_IMG, idx=10)
    closing_slide(prs, "Mensaje", ["línea"], "Acento naranja")
    prs.save("salida.pptx")

Requiere: pip install python-pptx
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Paleta oficial ----------
ORANGE = RGBColor(0xFD, 0x53, 0x1E)   # acento — resaltar lo importante
BLUE   = RGBColor(0x0A, 0x49, 0x9C)   # marca — SOLO logo, no diagramación
GRAY   = RGBColor(0x4A, 0x4A, 0x49)   # texto principal
DARK   = RGBColor(0x1A, 0x1C, 0x1D)   # títulos / casi negro
BANNER = RGBColor(0xF0, 0xF1, 0xF2)   # banda gris clara del título
ALT    = RGBColor(0xE4, 0xE4, 0xE6)   # filas alternas de tabla
SOFT   = RGBColor(0x8F, 0x98, 0x9D)   # gris medio (notas/pie)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Barlow"

EMU_IN = 914400
SW, SH = 13.333, 7.5
MARGIN = 0.7
CONTENT_W = SW - 2 * MARGIN
# Escala global de tipografía (una sola perilla → afecta a todos los cursos).
# 1.0 = tamaños base de diseño. (La "extensión" de una clase se controla con el
# NÚMERO de diapositivas de contenido por sesión, no con el tamaño de fuente.)
FONT_SCALE = 1.0
_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
LOGO = os.path.join(_ASSETS, "logo_nueva_america.png")
LOGO_WHITE = os.path.join(_ASSETS, "logo_nueva_america_blanco.png")   # para fondos oscuros/foto
GMAIL = os.path.join(_ASSETS, "gmail_icon.png")
FOTO_PORTADA = os.path.join(_ASSETS, "foto_portada_estudiantes.png")  # B/N estudiantes en sala
FOTO_ESTUDIANTE = os.path.join(_ASSETS, "foto_estudiante_tablet.jpg") # B/N estudiante con tablet
GESTORES_IMG = os.path.join(_ASSETS, "gestores_estudiantiles.png")    # flyer real de gestores
LINK = RGBColor(0x00, 0x97, 0xA7)     # teal del enlace de correo (como en el ejemplo)
_FOOTER = ""

def set_footer(text):
    global _FOOTER
    _FOOTER = text

# ---------- Texto enriquecido: **negrita** y @@naranja@@ ----------
def _run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    f = run.font
    f.size = Pt(round(size * FONT_SCALE, 1)); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color

def _rich(paragraph, markup, size, base_color, bold=False, italic=False):
    for part in re.split(r'(@@.*?@@|\*\*.*?\*\*)', markup):
        if part == "":
            continue
        run = paragraph.add_run()
        if part.startswith("@@") and part.endswith("@@"):
            _run(run, part[2:-2], size, ORANGE, bold=True, italic=italic)
        elif part.startswith("**") and part.endswith("**"):
            _run(run, part[2:-2], size, base_color, bold=True, italic=italic)
        else:
            _run(run, part, size, base_color, bold=bold, italic=italic)

# ---------- Primitivas ----------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg_white(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = WHITE

def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def add_logo(slide, width=2.7, mr=0.6, mb=0.45, corner="right"):
    if os.path.exists(LOGO):
        w = Inches(width)
        # alto según relación de aspecto del PNG (1498x314 ≈ 0.2096)
        h = Inches(width * 314.0 / 1498.0)
        left = Inches(mr) if corner == "left" else Emu(int(SW * EMU_IN) - w - Inches(mr))
        top = Emu(int(SH * EMU_IN) - h - Inches(mb))
        slide.shapes.add_picture(LOGO, left, top, width=w)

def _pic_fill(slide, img, x, y, w, h, alpha_pct=None):
    """Inserta una imagen RECORTADA para llenar la caja (x,y,w,h) sin deformarse.
    alpha_pct: opacidad 0-100 (p. ej. 19 = efecto 'lavado' de la portada institucional)."""
    if not os.path.exists(img):
        return None
    try:
        from PIL import Image as _PILImage
        iw, ih = _PILImage.open(img).size
    except Exception:
        iw, ih = (16, 9)
    box_ratio, img_ratio = w / h, iw / ih
    pic = slide.shapes.add_picture(img, Inches(x), Inches(y), Inches(w), Inches(h))
    # Recorte centrado (srcRect) para conservar la proporción de la foto
    if abs(img_ratio - box_ratio) > 1e-3:
        if img_ratio > box_ratio:   # imagen más ancha -> recortar lados
            crop = (1 - box_ratio / img_ratio) / 2
            pic.crop_left = crop; pic.crop_right = crop
        else:                       # imagen más alta -> recortar arriba/abajo
            crop = (1 - img_ratio / box_ratio) / 2
            pic.crop_top = crop; pic.crop_bottom = crop
    if alpha_pct is not None:
        blip = pic._element.blipFill.find(qn('a:blip'))
        mod = blip.makeelement(qn('a:alphaModFix'), {'amt': str(int(alpha_pct * 1000))})
        blip.append(mod)
    return pic

def veil(slide, x, y, w, h, color=BLUE, alpha_pct=45):
    """Velo de color semitransparente sobre una zona (duotono de foto, como en Objetivos)."""
    sp = rect(slide, x, y, w, h, color)
    srgb = sp._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    srgb.append(a)
    return sp

def footer_num(slide, idx):
    if idx is not None:
        tf = textbox(slide, SW - 1.2, SH - 0.45, 0.6, 0.3)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        _run(p.add_run(), str(idx), 10, SOFT, bold=True)

def title_block(slide, title, sub=None):
    """Banda gris clara con el título centrado (estilo base institucional)."""
    rect(slide, 0, 0, SW, 1.3, BANNER)
    tf = textbox(slide, MARGIN, 0, CONTENT_W, 1.3, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _rich(p, title, 30, DARK, bold=True)
    if sub:
        ts = textbox(slide, MARGIN, 1.4, CONTENT_W, 0.45)
        ps = ts.paragraphs[0]; ps.alignment = PP_ALIGN.CENTER
        _rich(ps, sub, 15, GRAY)
    return 2.05 if sub else 1.7

def bullets(slide, items, top, size=16, width=None, left=None):
    left = MARGIN if left is None else left
    width = CONTENT_W if width is None else width
    tf = textbox(slide, left, top, width, SH - top - 0.6)
    for i, it in enumerate(items):
        text, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(10 if lvl == 0 else 6)
        marker = "–   " if lvl == 0 else "●   "
        _rich(p, marker + text, size - (1 if lvl else 0), GRAY)
    return tf

# ---------- Tablas (sin bandas azules) ----------
def _nogrid(table):
    tblPr = table._tbl.tblPr
    for c in list(tblPr):
        if c.tag == qn('a:tableStyleId'):
            tblPr.remove(c)
    sid = tblPr.makeelement(qn('a:tableStyleId'), {})
    sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'  # No Style, No Grid
    tblPr.append(sid)
    table.first_row = False; table.horz_banding = False

def _fill(cell, color):
    cell.fill.solid(); cell.fill.fore_color.rgb = color

def _cell(cell, markup, size=12, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    _rich(p, markup, size, color, bold=bold)

def table_raw(slide, headers, rows, top, col_w=None, header_color=GRAY,
              fs_head=13, fs_body=12, aligns=None):
    ncol = len(headers); nrow = len(rows) + 1
    col_w = col_w or [CONTENT_W / ncol] * ncol
    s = sum(col_w); col_w = [w / s * CONTENT_W for w in col_w]
    height = min(SH - top - 0.6, 0.55 * nrow)
    table = slide.shapes.add_table(nrow, ncol, Inches(MARGIN), Inches(top),
                                   Inches(CONTENT_W), Inches(height)).table
    _nogrid(table)
    for j, w in enumerate(col_w):
        table.columns[j].width = Inches(w)
    aligns = aligns or [PP_ALIGN.LEFT] * ncol
    for j, h in enumerate(headers):
        c = table.cell(0, j); _fill(c, header_color)
        _cell(c, h, size=fs_head, color=WHITE, bold=True, align=aligns[j])
    for i, row in enumerate(rows):
        rc = WHITE if i % 2 == 0 else ALT
        for j, val in enumerate(row):
            c = table.cell(i + 1, j); _fill(c, rc)
            _cell(c, val, size=fs_body, color=GRAY, align=aligns[j])
    return table, top + height

# ---------- Slides de contenido ----------
def content_slide(prs, title, items, sub=None, size=16, idx=None, **_):
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    bullets(s, items, top=top + 0.15, size=size)
    footer_num(s, idx)
    return s

def table_content(prs, title, headers, rows, sub=None, note=None, col_w=None,
                  aligns=None, idx=None, fs_body=12, **_):
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    _, bottom = table_raw(s, headers, rows, top=top + 0.15, col_w=col_w, aligns=aligns, fs_body=fs_body)
    if note:
        _rich(textbox(s, MARGIN, min(SH - 0.85, bottom + 0.15), CONTENT_W, 0.6).paragraphs[0],
              note, 12, GRAY, italic=True)
    footer_num(s, idx)
    return s

# ---------- Portadas ----------
def course_cover(prs, materia, leccion, welcome, program_lines, start_note=None,
                 photo=None):
    """Portada del CURSO, réplica del deck institucional real (Internetworking):
    foto B/N LAVADA (19%) a la derecha · '¡Bienvenidos estudiantes!' en AZUL sobre la foto ·
    materia debajo · datos del programa arriba-izq · nota de inicio abajo-izq · logo abajo-der."""
    s = blank(prs); bg_white(s)
    photo = FOTO_PORTADA if photo is None else photo
    px = 4.2  # inicio de la foto (≈ tercio derecho, como en el ejemplo)
    _pic_fill(s, photo, px, 0, SW - px, SH, alpha_pct=19)
    # Datos del programa (arriba-izquierda)
    tp = textbox(s, MARGIN, 0.55, px - MARGIN + 2.0, 2.6)
    for i, ln in enumerate(program_lines):
        p = tp.paragraphs[0] if i == 0 else tp.add_paragraph()
        p.space_after = Pt(3); _rich(p, ln, 15, DARK)
    # Bienvenida (AZUL, como el deck real) + materia, centradas sobre la foto
    tw = textbox(s, px + 0.6, 1.6, SW - px - 1.2, 1.0)
    tw.paragraphs[0].alignment = PP_ALIGN.CENTER
    _rich(tw.paragraphs[0], welcome, 36, BLUE, bold=True)
    tm = textbox(s, px + 0.6, 3.15, SW - px - 1.2, 1.6)
    tm.paragraphs[0].alignment = PP_ALIGN.CENTER
    _rich(tm.paragraphs[0], materia, 28, DARK, bold=True)
    if leccion:
        pl = tm.add_paragraph(); pl.space_before = Pt(6); pl.alignment = PP_ALIGN.CENTER
        _rich(pl, leccion, 18, GRAY)
    # Franja blanca inferior (como el deck real): nota de inicio + logo
    rect(s, 0, SH - 0.95, SW, 0.95, WHITE)
    if start_note:
        _rich(textbox(s, MARGIN, SH - 0.72, 7.0, 0.5).paragraphs[0], start_note, 18, GRAY)
    add_logo(s, width=2.9, mb=0.18)
    return s

def session_cover(prs, kicker, nivel, title, subtitle, welcome, meta_lines, photo=None):
    """Portada de SESIÓN: etiqueta grande SESIÓN N + nivel + tema + foto lavada + logo."""
    s = blank(prs); bg_white(s)
    photo = FOTO_ESTUDIANTE if photo is None else photo
    if photo:
        px = 7.4
        _pic_fill(s, photo, px, 0, SW - px, SH, alpha_pct=19)
    rect(s, 0, 0, 0.30, SH, ORANGE)  # barra lateral de acento
    _rich(textbox(s, MARGIN, 0.5, CONTENT_W, 1.0).paragraphs[0], kicker, 46, DARK, bold=True)
    if nivel:
        _run(textbox(s, MARGIN, 1.7, CONTENT_W, 0.4).paragraphs[0].add_run(), nivel, 15, GRAY, bold=True)
    tt = textbox(s, MARGIN, 2.35, CONTENT_W - 0.5, 1.5)
    _rich(tt.paragraphs[0], title, 32, DARK, bold=True)
    if subtitle:
        psb = tt.add_paragraph(); psb.space_before = Pt(4); _rich(psb, subtitle, 22, GRAY)
    _rich(textbox(s, MARGIN, 4.2, CONTENT_W, 0.7).paragraphs[0], welcome, 24, ORANGE, bold=True)
    tmeta = textbox(s, MARGIN, 5.1, 8.0, 1.2)
    for i, ln in enumerate(meta_lines or []):
        p = tmeta.paragraphs[0] if i == 0 else tmeta.add_paragraph()
        p.space_after = Pt(4); _rich(p, ln, 14, GRAY)
    add_logo(s, width=2.6)
    return s

def objectives_slide(prs, label, items, idx=None, photo=None, title_in_card=None):
    """Slide de OBJETIVOS, réplica del deck real: foto a sangre a la derecha con VELO AZUL,
    tarjeta blanca a la izquierda con viñetas ● y etiqueta blanca sobre la foto."""
    s = blank(prs); bg_white(s)
    photo = FOTO_ESTUDIANTE if photo is None else photo
    px = 6.15  # inicio de la foto (≈ mitad derecha, como en el ejemplo)
    _pic_fill(s, photo, px, 0, SW - px, SH)
    veil(s, px, 0, SW - px, SH, color=BLUE, alpha_pct=45)
    # Tarjeta blanca (se monta sobre el borde izquierdo de la foto)
    rect(s, 0.15, 0.8, 8.0, 5.9, WHITE)
    # Etiqueta sobre la foto (derecha)
    tl = textbox(s, 8.45, 3.0, SW - 8.45 - 0.4, 0.9)
    tl.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tl.paragraphs[0].add_run(), label, 26, WHITE, bold=False)
    # Contenido de la tarjeta
    ty = 1.15
    if title_in_card:
        _rich(textbox(s, 0.65, ty, 7.0, 0.6).paragraphs[0], title_in_card, 20, DARK, bold=True)
        ty += 0.85
    tf = textbox(s, 0.65, ty, 7.0, 6.6 - ty)
    for i, it in enumerate(items):
        text, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(16)
        _rich(p, ("●   " if lvl == 0 else "–   ") + text, 18, GRAY)
    add_logo(s, width=2.4, corner="left")
    footer_num(s, idx)
    return s

def image_slide(prs, title, image_path, idx=None, sub=None):
    """Slide con banda de título + IMAGEN centrada que llena el área de contenido
    (p. ej. el flyer real de GESTORES ESTUDIANTILES)."""
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    area_top = top + 0.1
    area_h = SH - area_top - 0.35
    if os.path.exists(image_path):
        try:
            from PIL import Image as _PILImage
            iw, ih = _PILImage.open(image_path).size
        except Exception:
            iw, ih = (4, 3)
        ratio = iw / ih
        w = min(CONTENT_W + 0.8, area_h * ratio)
        h = w / ratio
        if h > area_h:
            h = area_h; w = h * ratio
        s.shapes.add_picture(image_path, Inches((SW - w) / 2), Inches(area_top + (area_h - h) / 2),
                             Inches(w), Inches(h))
    footer_num(s, idx)
    return s

def image_text_slide(prs, title, items, image_path, sub=None, size=15, idx=None, layout="auto", **_):
    """Slide MIXTA: viñetas + IMAGEN (diagrama, mockup o foto). Elige el layout para que la
    imagen se vea GRANDE y legible:
      - 'below' (imágenes ANCHAS): viñetas arriba a todo el ancho + imagen grande a todo el ancho abajo.
      - 'right' (imágenes altas/cuadradas): viñetas a la izquierda + imagen en panel derecho ancho.
      - 'auto' (default): 'below' si la imagen es ancha (ratio ≥ 1.35), si no 'right'.
    Si la imagen no existe, degrada a slide de solo texto."""
    if not image_path or not os.path.exists(image_path):
        return content_slide(prs, title, items, sub=sub, size=size, idx=idx)
    s = blank(prs); bg_white(s)
    top = title_block(s, title, sub)
    area_top = top + 0.15
    try:
        from PIL import Image as _PILImage
        iw, ih = _PILImage.open(image_path).size
    except Exception:
        iw, ih = (4, 3)
    ratio = iw / float(ih)
    mode = layout if layout != "auto" else ("below" if ratio >= 1.35 else "right")

    if mode == "below":
        # Viñetas arriba (ancho completo) + imagen GRANDE a todo el ancho abajo.
        avail_h = SH - area_top - 0.45
        img_band_h = avail_h * 0.58
        text_band_h = avail_h - img_band_h - 0.20
        bullets(s, items, top=area_top, size=size, left=MARGIN, width=CONTENT_W)
        maxw, maxh = CONTENT_W, img_band_h
        w = maxw; h = w / ratio
        if h > maxh:
            h = maxh; w = h * ratio
        px = MARGIN + (CONTENT_W - w) / 2
        py = area_top + text_band_h + 0.20 + (maxh - h) / 2
    else:
        # Panel derecho ANCHO (imágenes altas/cuadradas como mockups de teléfono).
        text_w = 5.9
        bullets(s, items, top=area_top, size=size, left=MARGIN, width=text_w)
        panel_x = MARGIN + text_w + 0.35
        panel_w = SW - MARGIN - panel_x
        panel_h = SH - area_top - 0.45
        w = panel_w; h = w / ratio
        if h > panel_h:
            h = panel_h; w = h * ratio
        px = panel_x + (panel_w - w) / 2
        py = area_top + (panel_h - h) / 2
    s.shapes.add_picture(image_path, Inches(px), Inches(py), Inches(w), Inches(h))
    footer_num(s, idx)
    return s

def tutor_slide(prs, nombre, credenciales, correo, idx=None):
    """Diapositiva '¿Quién soy?' (tutor), replicando el ejemplo institucional:
    nombre como título en banda gris, credenciales debajo, correo con ícono Gmail, logo abajo-izq."""
    s = blank(prs); bg_white(s)
    title_block(s, nombre)
    # Credenciales (indentadas a la izquierda, como en el ejemplo)
    tx = 3.0
    tc = textbox(s, tx, 2.5, SW - tx - MARGIN, 2.2)
    for i, ln in enumerate(credenciales):
        p = tc.paragraphs[0] if i == 0 else tc.add_paragraph()
        p.space_after = Pt(14); _rich(p, ln, 18, GRAY)
    # Correo con ícono de Gmail
    ey = 4.7
    if os.path.exists(GMAIL):
        s.shapes.add_picture(GMAIL, Inches(tx), Inches(ey), height=Inches(0.30))
    te = textbox(s, tx + 0.5, ey - 0.04, SW - tx - 0.5 - MARGIN, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    _run(te.paragraphs[0].add_run(), correo, 16, LINK, bold=False)
    add_logo(s, width=2.4, corner="left")
    footer_num(s, idx)
    return s

def closing_slide(prs, big, lines, accent):
    s = blank(prs); bg_white(s)
    rect(s, 0, 0, SW, 0.18, ORANGE)
    tf = textbox(s, MARGIN, 1.9, CONTENT_W, 1.2, anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _rich(tf.paragraphs[0], big, 40, DARK, bold=True)
    tl = textbox(s, MARGIN, 3.3, CONTENT_W, 1.6)
    for i, ln in enumerate(lines):
        p = tl.paragraphs[0] if i == 0 else tl.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6); _rich(p, ln, 16, GRAY)
    ta = textbox(s, MARGIN, 5.0, CONTENT_W, 0.7); ta.paragraphs[0].alignment = PP_ALIGN.CENTER
    _rich(ta.paragraphs[0], accent, 22, ORANGE, bold=True)
    add_logo(s, width=2.7)
    return s

def new_prs():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN)); prs.slide_height = Emu(int(SH * EMU_IN))
    cp = prs.core_properties
    cp.author = "Nueva América"; cp.last_modified_by = "Nueva América"
    cp.title = ""; cp.comments = ""
    return prs
