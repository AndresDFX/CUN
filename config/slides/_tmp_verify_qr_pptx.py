# -*- coding: utf-8 -*-
"""Verifica que las 5 Presentaciones del Curso incluyan el asset QR."""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

files = [
    Path(r"G:\Mi unidad\Trabajos\Empleo\CUN\Cursos\Especializacion\PROYECTO I\Clases\Presentacion del Curso - Proyecto I (nuevo).pptx"),
    Path(r"G:\Mi unidad\Trabajos\Empleo\CUN\Cursos\Pregrado\INVESTIGACION CIENCIA Y TECNOLOGIA PARA ESCUELA DE INGENIERIAS\Clases\Presentacion del Curso - Investigacion Ciencia y Tecnologia.pptx"),
    Path(r"G:\Mi unidad\Trabajos\Empleo\CUN\Cursos\Pregrado\CREATIVIDAD Y PENSAMIENTO INNOVADOR PARA ESCUELA DE INGENIERIAS\Clases\Presentacion del Curso - Creatividad y Pensamiento Innovador.pptx"),
    Path(r"G:\Mi unidad\Trabajos\Empleo\CUN\Cursos\Pregrado\TRABAJO DE GRADO 2 - MODELOS DE INNOVACION INGENIERIA DE SISTEMAS\Clases\Presentacion del Curso - Trabajo de Grado 2.pptx"),
    Path(r"G:\Mi unidad\Trabajos\Empleo\CUN\Cursos\Pregrado\TRABAJO DE GRADO 3 - MODELOS DE INNOVACION INGENIERIA DE SISTEMAS\Clases\Presentacion del Curso - Trabajo de Grado 3.pptx"),
]

for f in files:
    prs = Presentation(str(f))
    has_pic = False
    has_title = False
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # any picture beyond logo possible; check nearby title
                pass
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_pic = True
        blob = "\n".join(texts)
        if "PRESÉNTATE" in blob.upper() or "ROMPEHIELOS" in blob.upper() or "Escanea el QR" in blob:
            has_title = True
            # count pictures on this slide
            pics = sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
            print(f"OK {f.name}: slide {i} rompehielos, pics={pics}, total_slides={len(prs.slides)}")
            break
    else:
        print(f"FAIL {f.name}: sin slide rompehielos (pics_any={has_pic})")
