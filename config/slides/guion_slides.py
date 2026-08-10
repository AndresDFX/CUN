# -*- coding: utf-8 -*-
"""Mapa de slides REAL para los guiones docentes (y limpieza de la plantilla vieja).

Problema que resuelve
---------------------
Los guiones de la S02 en adelante se escribieron contra una plantilla fija de **7 slides**
(`Portada / OBJETIVOS / CONTENIDO CLAVE / RECUERDA-ENFOQUE / ACTIVIDAD-TALLER /
PARA CONTINUAR / Cierre`). Los `.pptx` reales se generan desde
`config/slides/content/cun_*.json` y tienen entre **16 y 26** slides con títulos propios:
la tabla «🗺️ Slides de esta presentación» describía un deck inexistente y el GUION
LITERAL mandaba leer en voz alta un número de slide que en pantalla muestra otra cosa.

Las S01 (encuadre) NO tienen este problema: su tabla está escrita a mano contra el deck
real, así que este módulo **solo** se aplica a los guiones que usaban la plantilla.

API
---
`titulos_pptx(pptx)`      → título real de cada slide (el motor no usa placeholder de
                            título: el título es el párrafo con la fuente más grande).
`tabla_slides_md(...)`    → tabla markdown con esos títulos reales.
`limpiar_referencias(md)` → quita de la narración los números de slide de la plantilla,
                            dejando el nombre del momento (Portada, TALLER, Cierre…).
                            No inventa números: los retira.
"""
from __future__ import annotations

import os
import re
import unicodedata

try:
    from pptx import Presentation
except Exception:  # pragma: no cover — sin python-pptx el guion no se degrada
    Presentation = None  # type: ignore


def titulos_pptx(pptx_path: str) -> list[str] | None:
    """Título de cada slide (índice 0 = slide 1). `None` si el deck no existe / no se lee."""
    if Presentation is None or not pptx_path or not os.path.isfile(pptx_path):
        return None
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return None
    out: list[str] = []
    for slide in prs.slides:
        best, best_size = "", -1.0
        title_shape = slide.shapes.title
        if title_shape is not None:
            try:
                best = title_shape.text.split("\n")[0].strip()
            except Exception:
                best = ""
        if not best:
            for sh in slide.shapes:
                if not getattr(sh, "has_text_frame", False):
                    continue
                for p in sh.text_frame.paragraphs:
                    txt = "".join(r.text for r in p.runs).strip()
                    if not txt:
                        continue
                    sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
                    size = max(sizes) if sizes else 0.0
                    if size > best_size:
                        best, best_size = txt, size
        out.append(best)
    return out


def deck_path(course_folder: str, label: str) -> str:
    """Ruta del deck de la sesión a partir de la carpeta del curso y el label `Sesion NN - …`."""
    return os.path.join(course_folder, "Clases", label, "Presentacion.pptx")


def tabla_slides_md(titulos: list[str] | None, *, encabezado: str | None = None) -> str | None:
    """Tabla «🗺️ Slides de esta presentación» con las slides REALES del deck."""
    if not titulos:
        return None
    enc = encabezado or (
        "🗺️ **Slides de esta presentación** "
        f"(deck real: **{len(titulos)} slides** — no es el mapa del curso)"
    )
    filas = [enc, "", "| Slide | Título en el PPTX |", "| :---: | :--- |"]
    for i, t in enumerate(titulos, 1):
        filas.append(f"| **{i}** | {(t or '—').strip()} |")
    filas.append("")
    return "\n".join(filas)


# ── Limpieza de las referencias de la plantilla de 7 slides ────────────────────
# Formas exactas encontradas en los guiones generados (ver auditoría 2026-08-10):
#   «#### 3️⃣ Taller (~20 min) — Slide 5»          → se quita el sufijo
#   «**Slides:** 1 (Portada) → 2 (OBJETIVOS)»      → «**Momento del deck:** Portada → OBJETIVOS»
#   «> “**Slide 6 — PARA CONTINUAR.** …»           → «> “**PARA CONTINUAR.** …»
#   «> “**Slide 3.** Tres palabras…»               → «> “Tres palabras…»
_RE_FASE = re.compile(r"(?m)^(#### [^\n]*?)\s+—\s+Slides?\s+\d{1,2}(?:\s*[–-]\s*\d{1,2})?\s*$")
_RE_SLIDES_LINEA = re.compile(r"(?m)^\*\*Slides:\*\*\s*(.+)$")
_RE_NUM_PAREN = re.compile(r"\b\d{1,2}\s*\(([^)]+)\)")
_RE_ROTULO = re.compile(r"\*\*Slide\s+\d{1,2}\s*[—–-]\s*")
_RE_ROTULO_COMA = re.compile(r"\*\*Slide\s+\d{1,2},\s*")
_RE_SOLO = re.compile(r"\*\*Slide\s+\d{1,2}\.\*\*\s*")


def limpiar_referencias(md: str) -> tuple[str, int]:
    """Retira los números de slide heredados de la plantilla. → (md, nº de cambios)."""
    n = 0

    def _cnt(pat, repl, texto):
        nonlocal n
        nuevo, k = pat.subn(repl, texto)
        n += k
        return nuevo

    md = _cnt(_RE_FASE, r"\1", md)

    def _slides_linea(m: re.Match) -> str:
        nonlocal n
        cuerpo = _RE_NUM_PAREN.sub(lambda g: g.group(1), m.group(1))
        n += 1
        return f"**Momento del deck:** {cuerpo}"

    md = _RE_SLIDES_LINEA.sub(_slides_linea, md)
    md = _cnt(_RE_ROTULO, "**", md)
    md = _cnt(_RE_ROTULO_COMA, "**", md)
    md = _cnt(_RE_SOLO, "", md)
    return md, n


NOTA_MOMENTOS = (
    "> Los **momentos** del plan de clase (Portada, OBJETIVOS, CONTENIDO CLAVE, TALLER, "
    "PARA CONTINUAR, Cierre) son los del guion, no números de slide: el deck real tiene "
    "más slides y su orden está en la tabla de arriba."
)


# ── Renumerado exacto para los guiones de S01 (mapa escrito a mano) ────────────
# El motor parte los bloques largos en slides «(cont.)», así que un mapa escrito a mano
# con N filas puede corresponder a un deck de N+2. A partir de la fila 7 de P1, por
# ejemplo, todos los números del guion quedaban corridos. Aquí se alinean las filas del
# mapa con los títulos reales (en orden) y se renumera la narración con el mapa exacto.
def _tokens(s: str) -> set[str]:
    s = unicodedata.normalize("NFKD", str(s or ""))
    s = "".join(c for c in s if not unicodedata.combining(c))
    return {w for w in re.findall(r"[a-z0-9]+", s.lower()) if len(w) > 2}


def _similitud(a: str, b: str) -> float:
    ta, tb = _tokens(a), _tokens(b)
    if not ta or not tb:
        return 0.0
    return len(ta & tb) / min(len(ta), len(tb))


def alinear_mapa(filas: list[tuple[int, str]], titulos: list[str]) -> dict[int, int]:
    """`{nº del guion: nº real}` alineando en orden las filas del mapa con el deck.

    Alineación monótona: cada fila del mapa avanza al siguiente título real que le
    corresponda; los títulos reales sobrantes son las slides «(cont.)» insertadas.
    """
    mapa: dict[int, int] = {}
    j = 0
    for k, (num, texto) in enumerate(filas):
        restantes = len(filas) - k
        mejor, mejor_s = None, 0.0
        # Solo se puede saltar hacia adelante lo que sobre de deck.
        tope = min(len(titulos), len(titulos) - restantes + 1)
        for i in range(j, max(j + 1, tope)):
            s = _similitud(texto, titulos[i])
            if s > mejor_s:
                mejor, mejor_s = i, s
        if mejor is None or mejor_s < 0.34:
            mejor = j
        mapa[num] = mejor + 1
        j = mejor + 1
    return mapa


_RE_CUALQUIER_SLIDE = re.compile(r"(Slides?)\s+(\d{1,2})")


def renumerar_referencias(md: str, mapa: dict[int, int]) -> tuple[str, int]:
    """Reescribe «Slide N» con el número real según `mapa`. → (md, nº de cambios)."""
    n = 0

    def _rep(m: re.Match) -> str:
        nonlocal n
        viejo = int(m.group(2))
        nuevo = mapa.get(viejo)
        if not nuevo or nuevo == viejo:
            return m.group(0)
        n += 1
        return f"{m.group(1)} {nuevo}"

    # No tocar la tabla de slides (ya viene del deck real): solo el cuerpo.
    corte = md.find("🎯")
    if corte < 0:
        return _RE_CUALQUIER_SLIDE.subn(_rep, md)
    cabeza, cuerpo = md[:corte], md[corte:]
    cuerpo = _RE_CUALQUIER_SLIDE.sub(_rep, cuerpo)
    return cabeza + cuerpo, n


_RE_FILA_MAPA = re.compile(r"^\| \*\*(\d{1,2})\*\* \| ([^|]+?)\s*\|", re.M)
# Algunos mapas agrupan varias slides en una fila: `| **7–8** | Mapa del curso |`.
_RE_FILA_RANGO = re.compile(r"^\| \*\*(\d{1,2})\s*[–—-]\s*(\d{1,2})\*\* \|", re.M)


def filas_mapa(md: str) -> list[tuple[int, str]]:
    """Filas `| **N** | Título |` del mapa de slides escrito a mano."""
    return [(int(a), b.strip()) for a, b in _RE_FILA_MAPA.findall(md)]


def _numeros_en_rangos(md: str) -> set[int]:
    """Slides ya cubiertas por filas de rango (`| **7–8** | … |`)."""
    cubiertos: set[int] = set()
    for a, b in _RE_FILA_RANGO.findall(md):
        cubiertos.update(range(int(a), int(b) + 1))
    return cubiertos


def ajustar_mapa_manual(md: str, titulos: list[str] | None) -> tuple[str, int]:
    """Realinea con el deck real un guion que trae mapa de slides escrito a mano (S01).

    Los mapas de la S01 están curados a mano y traen una columna extra («Cuándo usarla» /
    «Fase») que vale la pena conservar, así que no se sustituyen por la tabla automática:
    se **renumeran** contra el deck real y se añaden las filas de las slides «(cont.)»
    que el motor insertó. Si nada se movió, devuelve el guion intacto.
    """
    if not titulos:
        return md, 0
    filas = filas_mapa(md)
    if not filas:
        return md, 0
    mapa = alinear_mapa(filas, titulos)
    md, n = renumerar_referencias(md, mapa)

    usados = set(mapa.values()) | _numeros_en_rangos(md)
    ncols = md[md.find("| **%d**" % filas[0][0]):].split("\n", 1)[0].count("|") - 1
    for viejo, nuevo in sorted(mapa.items(), reverse=True):
        if viejo == nuevo:
            continue
        md = md.replace(f"| **{viejo}** |", f"| **{nuevo}** |", 1)
        n += 1
    # Filas que faltaban: slides «(cont.)» insertadas por el motor al partir un bloque.
    faltan = [i for i in range(1, len(titulos) + 1) if i not in usados]
    # Ascendente: al insertar la fila i, la i+1 ya encuentra su ancla `| **i** |`.
    for i in sorted(faltan):
        anterior = f"| **{i - 1}** |"
        if anterior not in md:
            continue
        linea_ant = md[md.index(anterior):].split("\n", 1)[0]
        extra = " |" * max(0, ncols - 1)
        md = md.replace(linea_ant, f"{linea_ant}\n| **{i}** | {titulos[i - 1]}{extra}", 1)
        n += 1
    return md, n
